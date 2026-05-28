"""최종 통합 발표 PPTX 생성기.

구성:
  [1p]    NEW 표지 — 데이터기반 도시설계 기말 프로젝트 + 팀 3인 + 학과·학번
  [2~16p] 메인 템플릿 1~15장 (DSM→DEM Converter, 그대로 유지)
  [17~26p] 그늘막 프로젝트 보완 슬라이드 10장
            - 17 왜 DSM 기반인가 (DSM vs DEM + 친구 작업과 연계 + 교수님 추천)
            - 18 이전 방식 → 보완 흐름
            - 19 데이터 9종 + Score 7 피처
            - 20 공간 필터링 4단계 + CV-A/CV-B
            - 21 CV-DSM + Self-consistency 5회 (NEW)
            - 22 흑석동 정밀 분석 (NGII + 실측 18개)
            - 23 가중치 자동 튜닝 (Bayesian Opt, NEW)
            - 24 OSMnx + 시나리오 6 + 강건 입지 2곳
            - 25 흑석동 사각지대 TOP10 + 예산 최적 배치
            - 26 결론 + 한 줄 요약 + 다음 단계

디자인: 메인 14p Palette/폰트 그대로 따름
  - 배경 #FFFFFF, 본문 #37352F, 액센트 #2383E2
  - 패널 #F7F6F3 / 보더 #E9E9E7
  - GREEN #0F7B6C, AMBER #CB7B26, RED #E03E3E
  - 폰트 Noto Sans KR (본문) + Consolas (코드)
"""
from __future__ import annotations
import shutil
from pathlib import Path
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


ROOT = Path(__file__).resolve().parent.parent
MAIN = ROOT / "presentation" / "데이터 도시기반 설계 과제 발표.pptx"
OUT = ROOT / "presentation" / "데이터기반_도시설계_기말_통합발표.pptx"


# 슬라이드별 스피커 노트 (PPTX notes_slide 에 자동 삽입)
# 메인 1~15 (PPTX 2~16) 는 친구 작품이라 손대지 않음
SLIDE_NOTES = {
    # 1: 새 표지
    1: """[1p · 표지 · 60초]
- 안녕하세요. 인공지능융합대학원 한영재입니다.
- 본 프로젝트는 3인 팀 작업입니다:
  · 문치국 — GIS·측량·BIM 엔지니어 (NGII·DSM 자문)
  · 한영재 — 한화투자증권 채권트레이더 (그늘막 분석 전체)
  · 원우식 — SOCAR ML/AI Engineer (SAM3·DSM→DEM 파이프라인)
- 메인은 두 부분으로 구성됩니다:
  1) 2~16p: 친구들이 만든 DSM → DEM Converter (15장)
  2) 17p~: 같은 DSM 자산을 재활용한 그늘막 입지 분석 (10장)
- 학제간 협업의 결과물입니다.
""",
    # 17: 왜 DSM 기반인가
    17: """[17p · 왜 DSM 기반인가 · 60초]
핵심 메시지: 친구들 작품 (DSM→DEM Converter) 의 DSM 자산을 우리가 재활용

DEM vs DSM (가장 자주 받는 질문):
- DEM: 객체 제거된 순수 지표 → 침수 분석·토목 설계
- DSM: 건물·수목 포함 표면 → 그림자·일조 분석
- 그림자를 만드는 건 표면 객체의 꼭대기. DEM 으론 그림자가 안 생김
- nDSM = DSM − DEM = 객체 순높이 = ray-cast 입력

연계 흐름:
- 친구 도구: 정사영상 → SAM3 객체 추출 → DSM 에서 제거 → GDAL 보간 → DEM
- 우리 도구: DSM 원본을 그대로 받아 → ray-cast 그림자 시뮬레이션 → 입지 추천

교수님 13p 추천: Deep Umbra (GAN 일조 컴퓨테이션)
- 풀 GAN 학습 부담 → DSM + 4시점 ray-cast 로 단순 이식
- 핵심 정신(다시점 시간 누적) 은 유지
""",
    # 18: 보완 흐름
    18: """[18p · 보완 흐름 · 45초]
v1 (초기) → v2 (최종) 비교 9 항목:
- 건물 입력: 더미 19동 → 실측 30+3,775동 (SAM + NGII)
- 그림자: 단일 시점 → 4시점 ray-cast (Deep Umbra 영감)
- 안정성: 결정론적 1회 → Self-consistency 5회 평균 (NEW)
- 보행로: 더미 14라인 buffer 20m → OSMnx walkable 5m
- 건물 cut: 없음 → NGII 5m inset (건물 위 추천 방지)
- 결집지: 없음 → 교차로+횡단보도 50m 근접 (OSMnx)
- 가중치: 수동 → Bayesian Opt 자동 튜닝 (NEW, 실측 18개 라벨)
- 정책: TOP10 정렬 → 배낭 최적화 (예산+분산 제약)

전체 메시지: 데이터 품질 + 시뮬레이션 정확도 + 정책 의사결정 3개 차원 동시 강화
""",
    # 19: 데이터 9종 + Score 7 피처
    19: """[19p · 데이터 9종 + Score · 45초]
3 계층:
- 기본 5종 (현재 4종 더미, 실데이터 자동 전환 구조)
- CV 2종 (V-World 항공 + Mapillary 거리뷰)
- 외부 정밀 2종 (NGII 흑석동 + OSMnx 도로망)

Score 식:
0.18·pop + 0.18·lst + 0.18·vuln − 0.15·shade − 0.05·natural
+ 0.12·sv_deficit + 0.20·intersection
→ 가중치는 다음 다음 슬라이드(23p) 에서 BayesOpt 로 자동 튜닝됨

격자 3,672 → 4단계 필터 → 330 → 시나리오 6 × TOP10 → 강건 입지
""",
    # 20: 공간 필터링 + CV
    20: """[20p · 공간 필터링 + CV · 45초]
4단계 필터 (3,672 → 330, 9% 살아남음):
- 보행로 5m: OSMnx walkable edges
- 건물 컷: NGII 폴리곤 5m inset (건물 위 추천 차단)
- 결집지 50m: 교차로 + 횡단보도 근접

CV-A · Mobile-SAM (Meta 2023):
- V-World 항공 z=15, 48 타일 합성
- zero-shot, 40MB CPU 추론, 학습 0건
- 건물 30동 자동 추출

CV-B · SegFormer-b0 (NVIDIA CityScapes pretrained):
- Mapillary 거리뷰 1,302장 → 후보 19격자 매핑
- 19 classes 의미 분할 → 그늘 결핍 공식
- TOP10 평균 deficit 0.153
""",
    # 21: CV-DSM + Self-consistency (NEW)
    21: """[21p · CV-DSM + Self-consistency · 60초 — NEW]
강의 13p + 40p 동시 반영:

CV-DSM (Deep Umbra 영감):
- 흑석동 DSM − DEM = nDSM (객체 높이)
- 4시점 (10·12·14·16시) shadow union → 누적 비율
- 309 격자 평균 누적 0.157

Self-consistency 5회 (NEW, 강의 40p 권장):
- 단일 ray-cast 는 태양위치 가정에 의존 → 신뢰구간 없음
- ±15분 시각, ±5° 방위 변동 5회 → 평균 + std
- 평균 std 0.021 (낮음 = 안정적, 결과 신뢰)
- 흑석동 focus map 에 std 레이어 토글 가능

예상 질문: 왜 풀 GAN 안 썼나? → 학습 시간·데이터 부담. ray-cast 로 같은 메시지
""",
    # 22: 흑석동 정밀 분석
    22: """[22p · 흑석동 정밀 분석 · 60초]
NGII 실측 3,775동 + 실측 그늘막 18개 (고정 13 + 스마트 5) + TOP10 분류

TOP10 3 카테고리:
- ⭐사각지대 (>150m): TOP3·4·5·8 (4곳) — 정책 어필 1순위
- ●보강 (<60m): TOP1·2·7·10 (4곳) — 기존 강화
- ○중간 (60~150m): TOP6·9 (2곳)

최대 사각지대: TOP4 (37.49786, 126.96129) — 359.6m
→ 흑석동 신규 설치 시 1순위 우선 후보

정책 메시지: 흑석동은 그늘막 75% 영역 커버 → 다른 동 우선 / 흑석동 내부 추가는 사각지대 4곳
""",
    # 23: 가중치 자동 튜닝 (NEW)
    23: """[23p · 가중치 자동 튜닝 · 75초 — NEW, 강의 40p 추천]
가설 (사용자 정정 반영):
- "이미 설치된 곳은 score 가 낮아야" → loss = mean(score(실측 18개 근접)) minimize

Bayesian Optimization (skopt.gp_minimize, n=60):
- GP surrogate 로 7 피처 가중치 자동 탐색
- 강의 40p 권장의 "데이터 기반 검증" 정신 반영

학습 결과 (수동 → 학습):
- shade −0.150 → −0.295  (페널티 2배 강화!)
- natural −0.050 → −0.178 (3.6배)
- lst +0.180 → +0.318 (폭염 우선)
- vuln +0.180 → +0.061 (흑석동 그늘막이 취약지에 집중되어 있어서 약화)

검증: 실측 위치 평균 score 0.087 → −0.122 (Δ −0.21, 가설 부합)
→ 알고리즘이 데이터로부터 "기존 설치 = 이미 좋은 곳" 을 자동 학습
""",
    # 24: OSMnx + 시나리오 + 강건 입지
    24: """[24p · OSMnx + 시나리오 + 강건 입지 · 60초]
OSMnx (보행자 결집 지점):
- walkable highway edges (보행로 5m 필터 입력)
- 교차로 노드 (street_count ≥ 3)
- 횡단보도 306 (한국 OSM 은 footway=crossing 우세 → 3종 태그 통합)

시나리오 6: 기본·고령자·폭염·유동인구·보행환경·교차로
- 가중치만 바꿔 재스코어 = 정책 민감도 분석
- 유동인구 최고 0.648, 3곳 독점

⭐ 강건 입지 2곳 (6 시나리오 모두 공통):
- 37.4907, 126.9647 — 동작대로 사당-이수 축
- 37.4898, 126.9670 — 동작대로 사당역 인근
→ 예산 제약 시 최우선 설치 대상
""",
    # 25: 사각지대 + 예산 최적화
    25: """[25p · 사각지대 + 예산 최적화 · 45초]
사각지대 4곳 (흑석동 한정, p22 참고):
- TOP4 359.6m (최대), TOP3 322.9m, TOP5 234.8m, TOP8 153.8m

예산 4천만원 배낭 최적화 (PuLP CBC, Optimal):
- 단가 800만원/개 → 최대 5개
- 200m 공간 분산 제약 (29쌍 active)
- 선정 5: 이수역·이수북·사당동·사당남·노량진
- 총 score 1.852, 예산 100% 사용

핵심 차이: TOP10 단순 정렬 vs 배낭 최적화
- 단순 정렬: 사당-이수에 7곳 몰림
- 배낭: 200m 분산으로 노량진역도 자동 포함
""",
    # 26: 결론 + 다음 단계
    26: """[26p · 결론 · 45초]
한 줄 요약:
"DSM → DEM Converter 자산을 재활용해 그늘막 입지를
 MCDA + CV + Self-consistency + Bayesian Opt + 배낭 최적화로 도출 →
 강건 입지 2 + 흑석동 사각지대 4 + 예산 4천만원 최적 배치 5"

다음 단계:
- 서울 열린데이터 API 연결 (P0)
- 동작구 전 동 실측 그늘막
- CV-D 멀티모달 VLM (강의 40~42p Reason-then-Estimate)
- 타 자치구 확장 (BBOX·NGII 교체)

감사합니다. 질문 받겠습니다.

예상 질문:
- 타 자치구 적용? → BBOX·NGII 교체. 가중치는 BayesOpt 로 재튜닝
- 코드? → github.com/thgtot92/project_1 (MIT)
""",
    # 27: 첨부 자료
    27: """[27p · 첨부 자료 · 30초]
- output/heukseok_focus.html — 흑석동 인터랙티브 지도 (브라우저)
- 같은 폴더에 동봉 / 더블클릭으로 열림

지도 레이어 (토글 가능):
- OSMnx 보행 도로 + 교차로 + 횡단보도
- NGII 실측 건물 3,775동 (층수별 색)
- 기존 그늘막 18개 (검정 우산)
- 수동 가중치 TOP10 (분홍 원)
- 학습 가중치 TOP10 (파란 별, BayesOpt 결과)
- 그림자 표준편차 (Self-consistency 5회)

심사위원이 직접 토글하면서 검토 가능
""",
}

# 메인 14p Palette (정확 추출)
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT     = RGBColor(0x37, 0x35, 0x2F)   # 본문 다크 그레이
C_SUB      = RGBColor(0x78, 0x77, 0x74)   # 서브 그레이
C_MUTED    = RGBColor(0x9B, 0x9A, 0x97)   # 더 흐린 그레이
C_ACCENT   = RGBColor(0x23, 0x83, 0xE2)   # 노션 블루
C_PANEL    = RGBColor(0xF7, 0xF6, 0xF3)   # 패널 베이지
C_BORDER   = RGBColor(0xE9, 0xE9, 0xE7)   # 보더
C_GREEN    = RGBColor(0x0F, 0x7B, 0x6C)
C_AMBER    = RGBColor(0xCB, 0x7B, 0x26)
C_RED      = RGBColor(0xE0, 0x3E, 0x3E)
FONT = "Noto Sans KR"
FONT_MONO = "Consolas"

# 슬라이드 크기 (10 × 5.625 inch)
TOTAL_NEW = 11    # 표지 1 + 보완 10
HEAD_LABEL_TOP   = Inches(0.47)
HEAD_TITLE_TOP   = Inches(0.72)
HEAD_DIVIDER_TOP = Inches(1.22)
BODY_TOP         = Inches(1.38)
FOOTER_TOP       = Inches(5.31)
LEFT             = Inches(0.68)
CONTENT_WIDTH    = Inches(8.65)
PAGE_NUM_LEFT    = Inches(8.57)


def _run(run, text, size=11, bold=False, color=C_TEXT, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _txt(slide, left, top, w, h, text, size=11, bold=False,
         color=C_TEXT, align=PP_ALIGN.LEFT, font=FONT):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = align
    _run(p.add_run(), text, size=size, bold=bold, color=color, font=font)
    return tx


def _bullets(slide, left, top, w, h, items, size=10, bullet="•  "):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p.add_run(), bullet + it, size=size, color=C_TEXT)
        p.space_after = Pt(4)
    return tx


def _box(slide, left, top, w, h, fill=C_PANEL, border=C_BORDER, border_w=0.75):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.adjustments[0] = 0.08  # 라운드 코너
    box.fill.solid(); box.fill.fore_color.rgb = fill
    if border is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = border
        box.line.width = Pt(border_w)
    box.shadow.inherit = False
    return box


def _divider(slide, top=HEAD_DIVIDER_TOP):
    """헤더 아래 1px 보더 라인."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   LEFT, top, CONTENT_WIDTH, Emu(9525))
    line.fill.solid(); line.fill.fore_color.rgb = C_BORDER
    line.line.fill.background()


def _header(slide, label, title):
    """메인 PPT 와 동일한 헤더 (작은 라벨 + 큰 제목 + 구분선)."""
    _txt(slide, LEFT, HEAD_LABEL_TOP, CONTENT_WIDTH, Inches(0.24),
         label, size=10, color=C_SUB)
    _txt(slide, LEFT, HEAD_TITLE_TOP, CONTENT_WIDTH, Inches(0.45),
         title, size=20, bold=True, color=C_TEXT)
    _divider(slide)


def _footer(slide, page):
    _txt(slide, LEFT, FOOTER_TOP, CONTENT_WIDTH, Inches(0.23),
         "동작구 그늘막 입지 분석 · 데이터기반 도시설계",
         size=9, color=C_MUTED)
    _txt(slide, PAGE_NUM_LEFT, FOOTER_TOP, Inches(0.75), Inches(0.23),
         f"{page:02d}", size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────
# 슬라이드 빌더들
# ─────────────────────────────────────────────────────────
def build_cover(slide):
    """[1p] 새 표지 — 제목 + 팀 + 학과·학번."""
    # 좌측 액센트 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.0), Inches(0.0),
                                   Inches(0.12), Inches(5.625))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    _txt(slide, Inches(0.68), Inches(0.7), Inches(9), Inches(0.3),
         "데이터기반 도시설계 · 기말 프로젝트", size=11, color=C_SUB)
    _txt(slide, Inches(0.68), Inches(1.1), Inches(9), Inches(0.7),
         "DSM → DEM Converter", size=24, bold=True, color=C_TEXT)
    _txt(slide, Inches(0.68), Inches(1.7), Inches(9), Inches(0.7),
         "& 동작구 여름 그늘막 최적 입지 추천", size=24, bold=True, color=C_ACCENT)
    # 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.68), Inches(2.5),
                                    Inches(2.0), Emu(38100))
    line.fill.solid(); line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()

    _txt(slide, Inches(0.68), Inches(2.75), Inches(9), Inches(0.3),
         "팀 구성 (3인)", size=11, bold=True, color=C_SUB)

    # 팀 카드 3개
    members = [
        ("문치국", "GIS · 측량 · BIM 엔지니어",
         "NGII 데이터 · DSM/DEM · 측량 도메인 자문"),
        ("한영재", "한화투자증권 · 채권 트레이더 / 인공지능융합대학원",
         "그늘막 입지 분석 (MCDA + CV + GIS + 배낭 최적화) · 전체 시스템"),
        ("원우식", "SOCAR · ML/AI Engineer",
         "SAM3 · GDAL 파이프라인 · DSM→DEM Converter"),
    ]
    card_w = Inches(2.95); card_h = Inches(1.4)
    for i, (name, role, contrib) in enumerate(members):
        x = Inches(0.68 + i * 3.05)
        c = _box(slide, x, Inches(3.15), card_w, card_h,
                  fill=C_PANEL, border=C_BORDER, border_w=0.75)
        tf = c.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        _run(p.add_run(), name, size=15, bold=True, color=C_TEXT)
        p2 = tf.add_paragraph()
        _run(p2.add_run(), role, size=9, color=C_ACCENT)
        p2.space_before = Pt(3)
        p3 = tf.add_paragraph()
        _run(p3.add_run(), contrib, size=9, color=C_SUB)
        p3.space_before = Pt(6)

    _txt(slide, Inches(0.68), Inches(4.85), Inches(9), Inches(0.3),
         "발표자 : 한영재 / 인공지능융합대학원 인공지능컴퓨팅 / 2025961227",
         size=10, color=C_TEXT)
    _txt(slide, Inches(0.68), Inches(5.15), Inches(9), Inches(0.3),
         "2026년 5월 · 연세대 도시공학과 · 데이터기반 도시설계",
         size=9, color=C_MUTED)


def build_p17_dsm_rationale(slide):
    """[17p] 왜 DSM 기반 그늘막 분석인가."""
    _header(slide, "연계 · 선정 이유",
            "왜 DSM 기반 그늘막 입지 분석인가")

    # 상단 부제
    _txt(slide, LEFT, BODY_TOP, CONTENT_WIDTH, Inches(0.4),
         "앞서 소개한 DSM→DEM Converter 의 자산을 그늘막 입지 분석에 재활용 — "
         "같은 데이터에서 다른 가치를 뽑아낸 학제간 협업.",
         size=10, color=C_SUB)

    # 좌측: DEM vs DSM
    _box(slide, LEFT, Inches(1.85), Inches(4.2), Inches(2.5),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(0.85), Inches(1.95), Inches(4), Inches(0.3),
         "DEM vs DSM", size=12, bold=True, color=C_ACCENT)
    _txt(slide, Inches(0.85), Inches(2.30), Inches(4), Inches(0.3),
         "DEM — 객체 제거된 순수 지표 (지형 분석)",
         size=9, color=C_TEXT)
    _txt(slide, Inches(0.85), Inches(2.55), Inches(4), Inches(0.3),
         "DSM — 건물·수목 포함 표면 (그림자 분석)",
         size=9, color=C_TEXT)
    _bullets(slide, Inches(0.85), Inches(2.90), Inches(4), Inches(1.5), [
        "그림자는 표면 객체(건물·수목)의 꼭대기가 광선을 가려야 생김",
        "DEM은 객체가 없으므로 그림자 시뮬레이션 불가",
        "DSM의 표면 = ray-cast 광선이 부딪힐 차폐물",
        "nDSM = DSM − DEM = 객체 순높이 (그림자 입력)",
    ], size=8)

    # 우측: 친구 작업과 우리 작업 연결
    _box(slide, Inches(5.05), Inches(1.85), Inches(4.28), Inches(2.5),
          fill=C_PANEL, border=C_ACCENT, border_w=1.25)
    _txt(slide, Inches(5.22), Inches(1.95), Inches(4), Inches(0.3),
         "친구 도구 → 우리 도구 자산 흐름", size=12, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(5.22), Inches(2.30), Inches(4), Inches(2.0), [
        "DSM(입력) — 그늘막 분석의 자연그늘 입력으로 재활용",
        "SAM3 건물 마스크 — 우리 CV-A SAM 의 검증 자료",
        "GDAL 보간 — CV-DSM ray-cast 결측치 보완에 응용",
        "DEM(출력) — 지형 분석 전용 (그늘막엔 사용 X)",
    ], size=8)

    # 하단: 교수님 13p 추천 — Deep Umbra
    _box(slide, LEFT, Inches(4.45), CONTENT_WIDTH, Inches(0.75),
          fill=RGBColor(0xFF, 0xF8, 0xE6), border=C_AMBER, border_w=0.75)
    _txt(slide, Inches(0.85), Inches(4.55), Inches(8.3), Inches(0.3),
         "📘 교수님 강의자료 13p 추천 — Deep Umbra (GAN 기반 일조 컴퓨테이션)",
         size=10, bold=True, color=C_AMBER)
    _txt(slide, Inches(0.85), Inches(4.85), Inches(8.3), Inches(0.3),
         "→ 본 프로젝트는 풀 GAN 학습 대신 DSM + 4시점 ray-cast 로 단순화 (핵심 정신: 다시점 시간 누적)",
         size=9, color=C_TEXT)

    _footer(slide, 17)


def build_p18_improvement(slide):
    """[18p] 이전 방식 → 보완 흐름."""
    _header(slide, "보완 흐름",
            "이전 방식 → DSM·교수님 추천·AI 튜닝으로 보완")

    # 표
    rows = [
        ("항목",            "이전 (v1)",         "보완 (v2, 최종)"),
        ("건물 입력",       "더미 19동 + 면적 추정 높이", "SAM 30동 (CV-A) + NGII 실측 3,775동 (흑석동)"),
        ("그림자 시뮬레이션", "오후 3시 단일 시점 평행이동",   "DSM 4시점 ray-cast (Deep Umbra 영감)"),
        ("그림자 안정성",    "결정론적 1회",        "Self-consistency 5회 평균 + 신뢰구간 (NEW)"),
        ("보행로 필터",     "더미 14개 라인 buffer 20m",   "OSMnx walkable highway buffer 5m"),
        ("건물 안 cut",     "없음",            "NGII 폴리곤 5m inset cut"),
        ("결집 지점",        "없음",            "OSMnx 교차로 + 횡단보도 50m 근접"),
        ("가중치",          "수동 7 피처",        "Bayesian Opt 자동 튜닝 (실측 18개 라벨, NEW)"),
        ("정책 의사결정",     "TOP10 정렬만",       "예산 + 200m 분산 → 배낭 최적화"),
    ]
    n_rows = len(rows)
    cw = [Inches(1.6), Inches(3.4), Inches(3.65)]
    tbl = slide.shapes.add_table(rows=n_rows, cols=3,
                                    left=LEFT, top=Inches(1.45),
                                    width=sum(cw, Inches(0)),
                                    height=Inches(3.5)).table
    for i, w in enumerate(cw):
        tbl.columns[i].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.margin_left = Inches(0.08)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            if ri == 0:
                _run(p.add_run(), val, size=10, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))
                cell.fill.solid(); cell.fill.fore_color.rgb = C_TEXT
            else:
                is_new = "(NEW)" in val
                color = C_ACCENT if (ci == 2 and is_new) else C_TEXT
                _run(p.add_run(), val, size=9, color=color,
                      bold=(ci == 0))
                if ri % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = C_PANEL

    _txt(slide, LEFT, Inches(5.05), CONTENT_WIDTH, Inches(0.3),
         "→ 데이터 품질·시뮬레이션 정확도·정책 의사결정 3개 차원 동시 강화",
         size=10, bold=True, color=C_ACCENT)
    _footer(slide, 18)


def build_p19_data_score(slide):
    """[19p] 데이터 9종 + Score 7 피처."""
    _header(slide, "입력 데이터 · Score",
            "9종 데이터를 7 피처로 통합 가중합")

    # 좌측: 데이터 9종 (압축)
    _txt(slide, LEFT, Inches(1.45), Inches(4.5), Inches(0.3),
         "데이터 9종 (3 계층)", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, LEFT, Inches(1.75), Inches(4.5), Inches(2.5), [
        "기본 5종: 생활인구 / LST / 취약계층 / 기존그늘막 / 인도",
        "CV 2종: V-World 항공 → SAM 건물 30동 (CV-A)",
        "         Mapillary 거리뷰 1,302장 (CV-B)",
        "외부 정밀 2종: NGII 흑석동 (DSM/DEM/건물 3,775동)",
        "              OSMnx 도로망·교차로·횡단보도 306",
    ], size=9)

    # 우측: Score 식
    _box(slide, Inches(5.3), Inches(1.45), Inches(4.03), Inches(1.6),
          fill=C_TEXT, border=None)
    _txt(slide, Inches(5.4), Inches(1.55), Inches(3.85), Inches(0.3),
         "Score = Σ wᵢ × MinMax(Featureᵢ)",
         size=10, bold=True, color=RGBColor(0xFF, 0xEE, 0x58), font=FONT_MONO)
    _txt(slide, Inches(5.4), Inches(1.85), Inches(3.85), Inches(0.3),
         "0.18·pop + 0.18·lst + 0.18·vuln",
         size=9, color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_MONO)
    _txt(slide, Inches(5.4), Inches(2.10), Inches(3.85), Inches(0.3),
         "− 0.15·shade − 0.05·natural",
         size=9, color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_MONO)
    _txt(slide, Inches(5.4), Inches(2.35), Inches(3.85), Inches(0.3),
         "+ 0.12·sv_deficit + 0.20·intersection",
         size=9, color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_MONO)
    _txt(slide, Inches(5.4), Inches(2.65), Inches(3.85), Inches(0.3),
         "(p23: 학습 가중치는 자동 튜닝 — Bayesian Opt)",
         size=8, color=C_AMBER)

    # 하단: 격자·시나리오
    _bullets(slide, LEFT, Inches(4.35), CONTENT_WIDTH, Inches(0.8), [
        "동작구 100m 격자 3,672 cells (EPSG:5179) · 보행로 + 건물 + 결집지 4단계 필터 → 330",
        "시나리오 6 (기본·고령자·폭염·유동인구·보행환경·교차로) → 강건 입지 식별",
    ], size=9)
    _footer(slide, 19)


def build_p20_filter_cv(slide):
    """[20p] 공간 필터링 4단계 + CV-A/CV-B."""
    _header(slide, "공간 필터링 + 컴퓨터비전",
            "차도·건물·녹지 추천 방지 + 항공·거리뷰 통합 분석")

    # 필터 퍼널 5단계
    funnel = [("3,672", "전체", C_SUB),
              ("449", "보행로 5m", C_AMBER),
              ("434", "건물 컷", C_RED),
              ("330", "결집지 50m", C_GREEN),
              ("10", "TOP", C_ACCENT)]
    for i, (n, l, c) in enumerate(funnel):
        x = LEFT + Inches(i * (1.75 + 0.05))
        box = _box(slide, x, Inches(1.45), Inches(1.75), Inches(1.1),
                    fill=c, border=None)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p.add_run(), n, size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        _run(p2.add_run(), l, size=9, color=RGBColor(0xFF, 0xFF, 0xFF))

    # CV-A
    _box(slide, LEFT, Inches(2.85), Inches(4.25), Inches(2.4),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(0.85), Inches(2.95), Inches(4), Inches(0.3),
         "CV-A · 항공 + Mobile-SAM", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(0.85), Inches(3.25), Inches(4), Inches(2.0), [
        "V-World WMTS z=15 · 48 타일 합성 (2048×1536)",
        "Mobile-SAM zero-shot (Meta 2023, 40MB CPU)",
        "면적·종횡비·밝기 필터 → 건물 30동",
        "흑석동에선 NGII 3,775동이 자동 union",
    ], size=8)

    # CV-B
    _box(slide, Inches(5.08), Inches(2.85), Inches(4.25), Inches(2.4),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(5.25), Inches(2.95), Inches(4), Inches(0.3),
         "CV-B · 거리뷰 + SegFormer", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(5.25), Inches(3.25), Inches(4), Inches(2.0), [
        "Mapillary 동작구 BBOX 4×4 → 1,302장",
        "SegFormer-b0 (CityScapes 19 classes pretrained)",
        "deficit = (도로+보도) × (1 − 건물 − 식생)",
        "TOP10 평균 deficit 0.153",
    ], size=8)
    _footer(slide, 20)


def build_p21_dsm_consistency(slide):
    """[21p] CV-DSM + Self-consistency 5회 (NEW)."""
    _header(slide, "CV-DSM · Self-consistency",
            "DSM 4시점 ray-cast + ±변동 5회 신뢰구간 (강의 13p + 40p 동시 반영)")

    _bullets(slide, LEFT, Inches(1.45), Inches(5.0), Inches(2.0), [
        "강의 13p 추천 Deep Umbra (GAN) → DSM ray-cast 로 단순 이식",
        "흑석동 DSM(표면) − DEM(지표) = nDSM (객체 순높이)",
        "4시점 (10·12·14·16시) shadow union → 누적 그림자 비율",
        "309 격자 평균 누적 그림자 0.157 → natural 피처",
    ], size=9)

    # Self-consistency 박스 (NEW)
    _box(slide, Inches(5.7), Inches(1.45), Inches(3.63), Inches(3.5),
          fill=C_PANEL, border=C_ACCENT, border_w=1.25)
    _txt(slide, Inches(5.85), Inches(1.55), Inches(3.4), Inches(0.3),
         "Self-consistency 5회 (NEW)", size=11, bold=True, color=C_ACCENT)
    _txt(slide, Inches(5.85), Inches(1.85), Inches(3.4), Inches(0.3),
         "강의 40p 권장 — 5회 반복 평균",
         size=9, color=C_SUB)
    _bullets(slide, Inches(5.85), Inches(2.15), Inches(3.4), Inches(2.8), [
        "태양 위치 ±15분 시각, ±5° 방위",
        "→ 5회 ray-cast 결과 평균 + std",
        "신뢰구간 [mean−2σ, mean+2σ] 산출",
        "안정적 격자 vs 변동 큰 격자 구분",
        "평균 std 0.021 (낮음 = 안정적)",
        "흑석동 309 격자에 적용",
    ], size=8)

    # 산출물
    _box(slide, LEFT, Inches(3.55), Inches(5.0), Inches(1.6),
          fill=RGBColor(0xF0, 0xF8, 0xFF), border=C_ACCENT, border_w=0.75)
    _txt(slide, Inches(0.85), Inches(3.65), Inches(4.6), Inches(0.3),
         "📁 산출물", size=10, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(0.85), Inches(3.95), Inches(4.6), Inches(1.2), [
        "heukseok_shadow_accum.tif (누적 그림자 래스터)",
        "grid_heukseok_consistency.csv (격자별 mean/std/ci95)",
        "흑석동 focus map v8 — 그림자 std 레이어 토글",
    ], size=8)
    _footer(slide, 21)


def build_p22_heukseok_focus(slide):
    """[22p] 흑석동 정밀 분석."""
    _header(slide, "흑석동 정밀 분석",
            "NGII 실측 건물 3,775동 + 실측 그늘막 18개 + 신규 추천 TOP 10")

    rows = [
        ("Rank", "score", "기존그늘막", "분류"),
        ("TOP1", "+0.276", "45.7m", "● 보강"),
        ("TOP2", "+0.231", "44.4m", "● 보강"),
        ("TOP3", "+0.134", "322.9m", "⭐ 사각지대"),
        ("TOP4", "+0.107", "359.6m", "⭐ 사각지대 (최대)"),
        ("TOP5", "+0.100", "234.8m", "⭐ 사각지대"),
        ("TOP6", "+0.072", "98.5m", "○ 중간"),
        ("TOP7", "+0.068", "51.0m", "● 보강"),
        ("TOP8", "+0.057", "153.8m", "⭐ 사각지대"),
        ("TOP9", "+0.016", "70.0m", "○ 중간"),
        ("TOP10", "−0.008", "29.2m", "● 보강"),
    ]
    cw = [Inches(0.75), Inches(0.95), Inches(1.15), Inches(1.55)]
    tbl = slide.shapes.add_table(rows=len(rows), cols=4,
                                    left=LEFT, top=Inches(1.45),
                                    width=sum(cw, Inches(0)),
                                    height=Inches(3.5)).table
    for i, w in enumerate(cw):
        tbl.columns[i].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.margin_left = Inches(0.05)
            tf = cell.text_frame
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            if ri == 0:
                _run(p.add_run(), val, size=9, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))
                cell.fill.solid(); cell.fill.fore_color.rgb = C_TEXT
            else:
                is_blind = "사각지대" in val
                color = C_RED if is_blind else C_TEXT
                _run(p.add_run(), val, size=8, color=color,
                      bold=is_blind)
                if is_blind:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xE5)

    # 우측 사각지대 강조
    _box(slide, Inches(4.7), Inches(1.45), Inches(4.63), Inches(2.0),
          fill=RGBColor(0xFF, 0xF4, 0xE5), border=C_AMBER, border_w=1.25)
    _txt(slide, Inches(4.85), Inches(1.55), Inches(4.4), Inches(0.3),
         "⭐ 사각지대 4곳 — 정책 어필 1순위", size=11, bold=True, color=C_AMBER)
    _bullets(slide, Inches(4.85), Inches(1.85), Inches(4.4), Inches(1.6), [
        "TOP4 (37.49786, 126.96129) — 최대 359.6m",
        "TOP3 (37.49874, 126.95789) — 322.9m",
        "TOP5 (37.51138, 126.96120) — 234.8m",
        "TOP8 (37.50328, 126.96465) — 153.8m",
    ], size=8)

    _box(slide, Inches(4.7), Inches(3.6), Inches(4.63), Inches(1.5),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(4.85), Inches(3.70), Inches(4.4), Inches(0.3),
         "흑석동 데이터", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(4.85), Inches(4.00), Inches(4.4), Inches(1.0), [
        "실측 그늘막 18개 (고정형 13 + 스마트형 5)",
        "NGII 건물 3,775동 (층수·HEGT) + DSM 4시점 누적",
        "그늘막 75% 영역 커버 → 정책: 다른 동 우선",
    ], size=8)
    _footer(slide, 22)


def build_p23_weight_tuning(slide):
    """[23p] 가중치 자동 튜닝 (Bayesian Opt, NEW)."""
    _header(slide, "가중치 자동 튜닝",
            "Bayesian Optimization (n=60) + 실측 그늘막 18개 역-라벨 학습 (NEW)")

    # 가설 박스
    _box(slide, LEFT, Inches(1.45), CONTENT_WIDTH, Inches(0.75),
          fill=RGBColor(0xF0, 0xF8, 0xFF), border=C_ACCENT, border_w=0.75)
    _txt(slide, Inches(0.85), Inches(1.55), Inches(8.3), Inches(0.3),
         "가설: 이미 그늘막이 설치된 위치는 score 가 낮아야 한다 (이미 충분히 커버됨)",
         size=10, bold=True, color=C_ACCENT)
    _txt(slide, Inches(0.85), Inches(1.85), Inches(8.3), Inches(0.3),
         "→ loss = mean( score(실측 그늘막 18개 근접 격자) )  — minimize",
         size=9, color=C_TEXT, font=FONT_MONO)

    # 가중치 비교 표
    rows = [
        ("피처",                  "수동 (직관)",  "학습 (BayesOpt)",  "변화"),
        ("popdens",              "+0.180",  "+0.260",  "↑ 강조"),
        ("lst",                  "+0.180",  "+0.318",  "↑↑ 폭염 우선"),
        ("vuln",                 "+0.180",  "+0.061",  "↓ 약화"),
        ("shade (페널티)",         "−0.150",  "−0.295",  "↑↑ 강화 (가설 부합)"),
        ("natural (페널티)",       "−0.050",  "−0.178",  "↑↑ 강화 (가설 부합)"),
        ("streetview_deficit",  "+0.120",  "+0.000",  "↓ 무력화"),
        ("intersection_density", "+0.200",  "+0.255",  "↑ 강조"),
    ]
    cw = [Inches(2.4), Inches(1.7), Inches(1.85), Inches(2.7)]
    tbl = slide.shapes.add_table(rows=len(rows), cols=4,
                                    left=LEFT, top=Inches(2.35),
                                    width=sum(cw, Inches(0)),
                                    height=Inches(2.4)).table
    for i, w in enumerate(cw):
        tbl.columns[i].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.margin_left = Inches(0.06)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            if ci > 0:
                p.alignment = PP_ALIGN.CENTER
            if ri == 0:
                _run(p.add_run(), val, size=9, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))
                cell.fill.solid(); cell.fill.fore_color.rgb = C_TEXT
            else:
                is_strong = "↑↑" in val
                color = C_RED if is_strong else C_TEXT
                _run(p.add_run(), val, size=8,
                      bold=(ci == 0 or is_strong),
                      color=color,
                      font=FONT_MONO if ci in (1, 2) else FONT)
                if is_strong:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xE5)

    # 검증 결과
    _box(slide, LEFT, Inches(4.85), CONTENT_WIDTH, Inches(0.45),
          fill=C_TEXT, border=None)
    _txt(slide, Inches(0.85), Inches(4.92), Inches(8.3), Inches(0.3),
         "검증: 실측 위치 평균 score   0.087 (수동) → −0.122 (학습)   "
         "Δ −0.21 — 학습이 실측 위치를 덜 추천 = 가설 부합",
         size=9, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    _footer(slide, 23)


def build_p24_osmnx_scenarios(slide):
    """[24p] OSMnx + 시나리오 6 + 강건 입지 2."""
    _header(slide, "OSMnx · 시나리오 · 강건 입지",
            "도로망·교차로·횡단보도 통합 + 6 정책 시나리오 비교")

    # OSMnx 통계
    _box(slide, LEFT, Inches(1.45), Inches(4.2), Inches(1.8),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(0.85), Inches(1.55), Inches(4), Inches(0.3),
         "OSMnx (보행자 결집 지점)", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(0.85), Inches(1.85), Inches(4), Inches(1.4), [
        "walkable highway edges — 보행로 5m 필터",
        "교차로 노드 (street_count ≥ 3)",
        "횡단보도 306 (highway/footway/crossing 3종 통합)",
        "결집지 union → 80m 내 개수 (피처) + 50m 근접 (필터)",
    ], size=8)

    # 시나리오 6 표
    rows = [
        ("시나리오",        "강조",                "최고",  "유니크"),
        ("기본",            "균형 18/18/18",       "0.551", "0"),
        ("고령자",          "vuln 0.38",          "0.518", "0"),
        ("폭염",            "lst 0.38",           "0.607", "1"),
        ("유동인구",         "pop 0.38",           "0.648", "3"),
        ("보행환경",         "sv_deficit 0.32",     "0.589", "0"),
        ("교차로",           "intersection 0.40",   "0.474", "0"),
    ]
    cw = [Inches(1.0), Inches(1.5), Inches(0.75), Inches(0.75)]
    tbl = slide.shapes.add_table(rows=len(rows), cols=4,
                                    left=Inches(5.83), top=Inches(1.45),
                                    width=sum(cw, Inches(0)),
                                    height=Inches(1.8)).table
    for i, w in enumerate(cw):
        tbl.columns[i].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.margin_left = Inches(0.04)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            if ci > 0:
                p.alignment = PP_ALIGN.CENTER
            if ri == 0:
                _run(p.add_run(), val, size=8, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))
                cell.fill.solid(); cell.fill.fore_color.rgb = C_TEXT
            else:
                _run(p.add_run(), val, size=8, color=C_TEXT,
                      bold=(ci == 0))

    # 강건 입지 카드
    _box(slide, LEFT, Inches(3.5), CONTENT_WIDTH, Inches(1.7),
          fill=RGBColor(0xFF, 0xF0, 0xF0), border=C_RED, border_w=1.25)
    _txt(slide, Inches(0.85), Inches(3.60), Inches(8.3), Inches(0.3),
         "⭐ 강건 입지 — 6 시나리오 모두 공통 추천 (정책 관점 불변)",
         size=11, bold=True, color=C_RED)
    _txt(slide, Inches(0.85), Inches(3.95), Inches(8.3), Inches(0.4),
         "#1  37.4907, 126.9647   동작대로 사당-이수 축",
         size=14, bold=True, color=C_TEXT)
    _txt(slide, Inches(0.85), Inches(4.35), Inches(8.3), Inches(0.4),
         "#2  37.4898, 126.9670   동작대로 사당역 인근",
         size=14, bold=True, color=C_TEXT)
    _txt(slide, Inches(0.85), Inches(4.85), Inches(8.3), Inches(0.3),
         "→ 예산 제약 시 최우선 설치 대상",
         size=9, color=C_SUB)
    _footer(slide, 24)


def build_p25_blindspot_budget(slide):
    """[25p] 사각지대 + 예산 최적화."""
    _header(slide, "사각지대 + 예산 최적화",
            "흑석동 사각지대 4곳 + 예산 4천만원 배낭 최적 배치 5곳")

    # 좌측: 사각지대
    _box(slide, LEFT, Inches(1.45), Inches(4.2), Inches(2.7),
          fill=RGBColor(0xFF, 0xF4, 0xE5), border=C_AMBER, border_w=1.25)
    _txt(slide, Inches(0.85), Inches(1.55), Inches(4), Inches(0.3),
         "흑석동 사각지대 TOP 4", size=11, bold=True, color=C_AMBER)
    _bullets(slide, Inches(0.85), Inches(1.85), Inches(4), Inches(2.3), [
        "TOP4 (37.49786, 126.96129) — 359.6m  (최대)",
        "TOP3 (37.49874, 126.95789) — 322.9m",
        "TOP5 (37.51138, 126.96120) — 234.8m",
        "TOP8 (37.50328, 126.96465) — 153.8m",
        "기존 그늘막 18개와 비교 — 정책 어필 1순위",
    ], size=8)

    # 우측: 예산 최적화
    _box(slide, Inches(5.7), Inches(1.45), Inches(3.63), Inches(2.7),
          fill=C_PANEL, border=C_ACCENT, border_w=1.25)
    _txt(slide, Inches(5.85), Inches(1.55), Inches(3.4), Inches(0.3),
         "예산 4천만원 최적 배치", size=11, bold=True, color=C_ACCENT)
    _txt(slide, Inches(5.85), Inches(1.85), Inches(3.4), Inches(0.3),
         "PuLP 정수 선형계획법 (CBC, Optimal)",
         size=8, color=C_SUB)
    _bullets(slide, Inches(5.85), Inches(2.15), Inches(3.4), Inches(2.0), [
        "단가 800만원/개 → 최대 5개",
        "공간 분산 ≥ 200m (29쌍 제약)",
        "선정 5: 이수역·이수북·사당동·사당남·노량진",
        "총 score 1.852, 예산 100% 사용",
    ], size=8)

    # 하단: 핵심 메시지
    _box(slide, LEFT, Inches(4.30), CONTENT_WIDTH, Inches(0.75),
          fill=C_TEXT, border=None)
    _txt(slide, Inches(0.85), Inches(4.40), Inches(8.3), Inches(0.3),
         "TOP10 단순 정렬 vs 배낭 최적화 — 200m 분산 제약으로 같은 곳 몰림 방지",
         size=10, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    _txt(slide, Inches(0.85), Inches(4.70), Inches(8.3), Inches(0.3),
         "→ 노량진역이 자동 포함되며 정책 효용 ↑",
         size=9, color=RGBColor(0xFF, 0xFF, 0xFF))
    _footer(slide, 25)


def build_p27_attachment(slide):
    """[27p] 첨부 자료 — 흑석동 HTML 인터랙티브 지도 안내."""
    _header(slide, "첨부 자료",
            "흑석동 인터랙티브 지도 — output/heukseok_focus.html")

    _txt(slide, LEFT, Inches(1.45), CONTENT_WIDTH, Inches(0.4),
         "본 PPTX 와 함께 동봉된 HTML 파일을 브라우저에서 열면 "
         "심사위원이 직접 레이어를 토글하며 검토할 수 있습니다.",
         size=10, color=C_SUB)

    # 좌측: 파일 경로 박스
    _box(slide, LEFT, Inches(2.0), Inches(4.3), Inches(1.8),
          fill=C_TEXT, border=None)
    _txt(slide, Inches(0.85), Inches(2.15), Inches(4), Inches(0.3),
         "📎 첨부 파일", size=11, bold=True,
         color=RGBColor(0xFF, 0xEE, 0x58))
    _txt(slide, Inches(0.85), Inches(2.55), Inches(4), Inches(0.4),
         "output/heukseok_focus.html",
         size=12, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_MONO)
    _txt(slide, Inches(0.85), Inches(3.05), Inches(4), Inches(0.4),
         "약 7 MB · 단일 HTML (의존성 없음)",
         size=9, color=RGBColor(0xCC, 0xCC, 0xCC))
    _txt(slide, Inches(0.85), Inches(3.35), Inches(4), Inches(0.4),
         "더블클릭 → 기본 브라우저로 자동 열림",
         size=9, color=RGBColor(0xCC, 0xCC, 0xCC))

    # 우측: 지도 레이어 목록
    _box(slide, Inches(5.18), Inches(2.0), Inches(4.15), Inches(2.8),
          fill=C_PANEL, border=C_ACCENT, border_w=1.25)
    _txt(slide, Inches(5.33), Inches(2.10), Inches(4), Inches(0.3),
         "지도 레이어 (토글 가능)", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(5.33), Inches(2.40), Inches(4), Inches(2.4), [
        "OSMnx 보행 도로 + 교차로 + 횡단보도",
        "NGII 실측 건물 3,775동 (층수별 색)",
        "기존 그늘막 18개 (검정/진청 우산)",
        "수동 가중치 TOP10 (분홍 원)",
        "학습 가중치 TOP10 (파란 별, BayesOpt)",
        "그림자 표준편차 (Self-consistency 5회)",
    ], size=8)

    # 하단: 사용 안내
    _box(slide, LEFT, Inches(4.20), CONTENT_WIDTH, Inches(0.85),
          fill=RGBColor(0xF0, 0xF8, 0xFF), border=C_ACCENT, border_w=0.75)
    _txt(slide, Inches(0.85), Inches(4.30), Inches(8.3), Inches(0.3),
         "💡 발표 직후 시연 가능 — 브라우저 창에 미리 열어두기 권장",
         size=10, bold=True, color=C_ACCENT)
    _txt(slide, Inches(0.85), Inches(4.60), Inches(8.3), Inches(0.4),
         "사각지대(TOP4·3·5·8) 4곳을 지도에서 직접 확인 → 가장 가까운 기존 그늘막까지 거리 popup",
         size=9, color=C_TEXT)
    _footer(slide, 27)


def build_p26_conclusion(slide):
    """[26p] 결론 + 한 줄 요약 + 다음 단계."""
    _header(slide, "결론 · 다음 단계", "")

    # 좌측: 한계
    _txt(slide, LEFT, Inches(1.45), Inches(4.3), Inches(0.3),
         "현재 한계", size=11, bold=True, color=C_RED)
    _bullets(slide, LEFT, Inches(1.75), Inches(4.3), Inches(1.8), [
        "5종 기본 데이터 중 4종 더미",
        "LST 래스터 직접 파싱 미연결",
        "SAM 추출 건물 추정 높이 (NGII는 실측)",
        "예산 단가 800만원/개 가정 (실 단가 확보 시 교체)",
    ], size=9)

    # 우측: 다음 단계
    _txt(slide, Inches(5.05), Inches(1.45), Inches(4.3), Inches(0.3),
         "다음 단계", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(5.05), Inches(1.75), Inches(4.3), Inches(1.8), [
        "서울 열린데이터 API 연결 (P0)",
        "동작구 전 동 실측 그늘막 (흑석동만 확보)",
        "CV-D 멀티모달 VLM (강의 40~42p Reason-then-Estimate)",
        "타 자치구 확장 (BBOX·NGII 교체)",
    ], size=9)

    # 한 줄 요약
    _box(slide, LEFT, Inches(3.85), CONTENT_WIDTH, Inches(1.3),
          fill=C_TEXT, border=None)
    _txt(slide, Inches(0.85), Inches(3.95), Inches(8.3), Inches(0.4),
         "DSM → DEM Converter 의 자산 (DSM·SAM3·GDAL) 을 재활용해",
         size=12, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    _txt(slide, Inches(0.85), Inches(4.30), Inches(8.3), Inches(0.4),
         "그늘막 입지를 MCDA + CV + Self-consistency + Bayesian Opt + 배낭 최적화로 도출",
         size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _txt(slide, Inches(0.85), Inches(4.70), Inches(8.3), Inches(0.3),
         "→ 강건 입지 2 + 흑석동 사각지대 4 + 예산 4천만원 최적 배치 5",
         size=10, color=RGBColor(0xCC, 0xCC, 0xCC))
    _footer(slide, 26)


# ─────────────────────────────────────────────────────────
# 메인 빌더
# ─────────────────────────────────────────────────────────
def move_slide_to_front(prs, src_idx):
    """xml 직접 조작으로 src_idx 슬라이드를 맨 앞으로."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    target = slides[src_idx]
    xml_slides.remove(target)
    xml_slides.insert(0, target)


def build():
    # 1) 메인 PPT 복사
    shutil.copy(MAIN, OUT)
    prs = Presentation(str(OUT))
    blank = prs.slide_layouts[6]  # blank layout (메인과 동일)

    print(f"메인 슬라이드: {len(prs.slides)}장 (그대로 유지)")

    # 2) 끝에 새 표지 + 보완 10장 추가
    builders = [
        build_cover,
        build_p17_dsm_rationale,
        build_p18_improvement,
        build_p19_data_score,
        build_p20_filter_cv,
        build_p21_dsm_consistency,
        build_p22_heukseok_focus,
        build_p23_weight_tuning,
        build_p24_osmnx_scenarios,
        build_p25_blindspot_budget,
        build_p26_conclusion,
        build_p27_attachment,
    ]
    cover_idx = len(prs.slides)  # 추가될 표지의 위치 (15)
    for builder in builders:
        slide = prs.slides.add_slide(blank)
        # layout 에서 상속된 placeholder (예: "제목 입력" 박스) 모두 제거
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        # 흰 배경
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                       prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()
        # 가장 뒤로 보내기 (다른 요소들이 위에 그려지도록)
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)
        builder(slide)
        print(f"  + 추가 {len(prs.slides):2d}: {builder.__name__}")

    # 3) 표지(가장 마지막 추가 - 사실 첫 번째 builder가 cover)를 맨 앞으로
    # cover_idx 가 cover 슬라이드 위치임
    move_slide_to_front(prs, cover_idx)
    print(f"  → 표지를 1번째로 이동 (메인은 2~16, 보완은 17~27)")

    # 4) 스피커 노트 자동 삽입 (1p + 17~27p, 메인 2~16 은 친구 작품이라 건드리지 않음)
    notes_inserted = 0
    for i, slide in enumerate(prs.slides, 1):
        note = SLIDE_NOTES.get(i, "")
        if note:
            slide.notes_slide.notes_text_frame.text = note.strip()
            notes_inserted += 1
    print(f"  → 스피커 노트 {notes_inserted}장 (1p + 17~27p) 삽입 완료")

    prs.save(str(OUT))
    print(f"\n[완료] {OUT}")
    print(f"  총 {len(prs.slides)}장")


if __name__ == "__main__":
    build()
