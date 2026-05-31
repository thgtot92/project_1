"""중간 발표용 PPTX 8장 생성.

구성:
  1p 표지 — 중간 발표 · 팀 3인 · 학과·학번
  2p 프로젝트 개요 + 두 트랙 소개
  3p 팀 구성과 역할 분담 (3명 상세)
  4p 일정·마일스톤
  5p 트랙 1: DSM→DEM Converter — 진행 상황·기술 스택
  6p 트랙 2: 그늘막 입지 분석 — 진행 상황·기술 스택
  7p 두 트랙 연계와 학제간 협업
  8p 남은 작업 + 최종 발표 준비

디자인: 메인 PPT 14p Palette 그대로 (Noto Sans KR, 노션 스타일)
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "presentation" / "중간발표_데이터기반_도시설계.pptx"

# Palette (메인 14p와 동일)
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT     = RGBColor(0x37, 0x35, 0x2F)
C_SUB      = RGBColor(0x78, 0x77, 0x74)
C_MUTED    = RGBColor(0x9B, 0x9A, 0x97)
C_ACCENT   = RGBColor(0x23, 0x83, 0xE2)
C_PANEL    = RGBColor(0xF7, 0xF6, 0xF3)
C_BORDER   = RGBColor(0xE9, 0xE9, 0xE7)
C_GREEN    = RGBColor(0x0F, 0x7B, 0x6C)
C_AMBER    = RGBColor(0xCB, 0x7B, 0x26)
C_RED      = RGBColor(0xE0, 0x3E, 0x3E)
FONT = "Noto Sans KR"
FONT_MONO = "Consolas"

TOTAL = 8
LEFT             = Inches(0.68)
CONTENT_WIDTH    = Inches(8.65)
HEAD_LABEL_TOP   = Inches(0.47)
HEAD_TITLE_TOP   = Inches(0.72)
HEAD_DIVIDER_TOP = Inches(1.22)
BODY_TOP         = Inches(1.38)
FOOTER_TOP       = Inches(5.31)
PAGE_NUM_LEFT    = Inches(8.57)


# ─────────── 헬퍼 ───────────
def _run(run, text, size=11, bold=False, color=C_TEXT, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _txt(slide, left, top, w, h, text, size=11, bold=False,
         color=C_TEXT, align=PP_ALIGN.LEFT, font=FONT):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = align
    _run(p.add_run(), text, size=size, bold=bold, color=color, font=font)
    return tx


def _bullets(slide, left, top, w, h, items, size=10, bullet="•  "):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p.add_run(), bullet + it, size=size, color=C_TEXT)
        p.space_after = Pt(4)
    return tx


def _box(slide, left, top, w, h, fill=C_PANEL, border=C_BORDER, border_w=0.75):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.adjustments[0] = 0.08
    box.fill.solid(); box.fill.fore_color.rgb = fill
    if border is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = border
        box.line.width = Pt(border_w)
    box.shadow.inherit = False
    return box


def _divider(slide, top=HEAD_DIVIDER_TOP):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   LEFT, top, CONTENT_WIDTH, Emu(9525))
    line.fill.solid(); line.fill.fore_color.rgb = C_BORDER
    line.line.fill.background()


def _header(slide, label, title):
    _txt(slide, LEFT, HEAD_LABEL_TOP, CONTENT_WIDTH, Inches(0.24),
         label, size=10, color=C_SUB)
    _txt(slide, LEFT, HEAD_TITLE_TOP, CONTENT_WIDTH, Inches(0.45),
         title, size=20, bold=True, color=C_TEXT)
    _divider(slide)


def _footer(slide, page):
    _txt(slide, LEFT, FOOTER_TOP, CONTENT_WIDTH, Inches(0.23),
         "중간 발표 · 데이터기반 도시설계 · 한영재·문치국·원우식",
         size=9, color=C_MUTED)
    _txt(slide, PAGE_NUM_LEFT, FOOTER_TOP, Inches(0.75), Inches(0.23),
         f"{page:02d} / {TOTAL:02d}", size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # placeholder 제거
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    # 흰 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


# ─────────── 슬라이드 빌더 ───────────
def build_cover(slide):
    """[1p] 중간 발표 표지."""
    # 좌측 액센트 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.0), Inches(0.0),
                                   Inches(0.12), Inches(5.625))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    _txt(slide, Inches(0.68), Inches(0.7), Inches(9), Inches(0.3),
         "데이터기반 도시설계 · 중간 발표", size=11, color=C_SUB)
    _txt(slide, Inches(0.68), Inches(1.1), Inches(9), Inches(0.7),
         "DSM → DEM Converter", size=22, bold=True, color=C_TEXT)
    _txt(slide, Inches(0.68), Inches(1.65), Inches(9), Inches(0.7),
         "& 동작구 여름 그늘막 입지 분석", size=22, bold=True, color=C_ACCENT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.68), Inches(2.45),
                                    Inches(2.0), Emu(38100))
    line.fill.solid(); line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()

    # 진행 상태 박스
    _box(slide, Inches(0.68), Inches(2.70), Inches(8.65), Inches(0.55),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(0.85), Inches(2.83), Inches(8.4), Inches(0.3),
         "중간 점검 발표 (2026-06-04) — 진행 상황·역할 분담·향후 계획 공유",
         size=11, color=C_TEXT)

    # 팀 카드 3개
    members = [
        ("문치국", "GIS · 측량 · BIM 엔지니어",
         "NGII 데이터 · DSM/DEM · 측량 도메인 자문"),
        ("한영재", "한화투자증권 · 채권 트레이더\n인공지능융합대학원",
         "그늘막 입지 분석 · 전체 시스템 통합"),
        ("원우식", "SOCAR · ML/AI Engineer",
         "SAM3 · GDAL · DSM→DEM Converter"),
    ]
    card_w = Inches(2.85); card_h = Inches(1.35)
    for i, (name, role, contrib) in enumerate(members):
        x = Inches(0.68 + i * 2.95)
        c = _box(slide, x, Inches(3.50), card_w, card_h,
                  fill=C_PANEL, border=C_BORDER, border_w=0.75)
        tf = c.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.13)
        p = tf.paragraphs[0]
        _run(p.add_run(), name, size=14, bold=True, color=C_TEXT)
        for line_txt in role.split("\n"):
            p = tf.add_paragraph()
            _run(p.add_run(), line_txt, size=8, color=C_ACCENT)
            p.space_before = Pt(2)
        p = tf.add_paragraph()
        _run(p.add_run(), contrib, size=8, color=C_SUB)
        p.space_before = Pt(5)

    _txt(slide, Inches(0.68), Inches(5.05), Inches(9), Inches(0.3),
         "발표자 : 한영재 / 인공지능융합대학원 인공지능컴퓨팅 / 2025961227",
         size=10, color=C_TEXT)
    _txt(slide, Inches(0.68), Inches(5.32), Inches(9), Inches(0.3),
         "2026년 6월 · 연세대 도시공학과 · 데이터기반 도시설계",
         size=9, color=C_MUTED)


def build_overview(slide):
    """[2p] 프로젝트 개요 + 두 트랙."""
    _header(slide, "개요", "두 개의 트랙으로 진행 중인 학제간 프로젝트")

    _txt(slide, LEFT, Inches(1.45), CONTENT_WIDTH, Inches(0.4),
         "동일한 DSM 자산을 두 가지 가치로 변환하는 학제간 협업 프로젝트.",
         size=10, color=C_SUB)

    # 두 트랙 박스
    # 트랙 1
    _box(slide, LEFT, Inches(1.95), Inches(4.15), Inches(2.7),
          fill=C_PANEL, border=C_BORDER, border_w=0.75)
    _txt(slide, Inches(0.85), Inches(2.05), Inches(4), Inches(0.3),
         "트랙 1", size=10, bold=True, color=C_SUB)
    _txt(slide, Inches(0.85), Inches(2.30), Inches(4), Inches(0.4),
         "DSM → DEM Converter", size=14, bold=True, color=C_ACCENT)
    _txt(slide, Inches(0.85), Inches(2.65), Inches(4), Inches(0.3),
         "정사영상에서 SAM3로 객체 추출 → GDAL 보간 → 순수 지형 DEM",
         size=9, color=C_SUB)
    _bullets(slide, Inches(0.85), Inches(3.00), Inches(3.9), Inches(1.6), [
        "담당: 문치국 + 원우식",
        "기술: PyQt6 · SAM3 · GDAL · rasterio",
        "산출: DEM GeoTIFF + 인터랙티브 지도",
        "용도: 침수·토목 설계 등 지형 분석",
    ], size=9)

    # 트랙 2
    _box(slide, Inches(5.18), Inches(1.95), Inches(4.15), Inches(2.7),
          fill=C_PANEL, border=C_ACCENT, border_w=1.25)
    _txt(slide, Inches(5.35), Inches(2.05), Inches(4), Inches(0.3),
         "트랙 2", size=10, bold=True, color=C_SUB)
    _txt(slide, Inches(5.35), Inches(2.30), Inches(4), Inches(0.4),
         "그늘막 입지 분석", size=14, bold=True, color=C_ACCENT)
    _txt(slide, Inches(5.35), Inches(2.65), Inches(4), Inches(0.3),
         "DSM 원본을 재활용해 동작구 여름 그늘막 최적 입지 추천",
         size=9, color=C_SUB)
    _bullets(slide, Inches(5.35), Inches(3.00), Inches(3.9), Inches(1.6), [
        "담당: 한영재 (메인)",
        "기술: MCDA · CV(SAM·SegFormer) · OSMnx · PuLP",
        "산출: TOP 추천 + 인터랙티브 지도",
        "용도: 정책 의사결정 (예산 4천만원 배치 등)",
    ], size=9)

    # 하단: 연결고리 강조
    _box(slide, LEFT, Inches(4.85), CONTENT_WIDTH, Inches(0.45),
          fill=C_TEXT, border=None)
    _txt(slide, Inches(0.85), Inches(4.92), Inches(8.3), Inches(0.3),
         "🔗 연결고리: 같은 DSM 자산을 두 도구가 다른 용도로 활용 — 학제간 협업의 핵심",
         size=10, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    _footer(slide, 2)


def build_team(slide):
    """[3p] 팀 구성과 역할 분담."""
    _header(slide, "팀 구성", "3인 학제간 협업 — 각자의 전문성을 모듈로 분리")

    # 표 형태로 팀원 + 역할 + 담당 모듈
    rows = [
        ("팀원",       "전문 분야",                     "본 프로젝트 담당"),
        ("문치국",     "GIS · 측량 · BIM 엔지니어",      "NGII 데이터 확보 (DSM·DEM·건물 shp)\n측량 도메인 자문 · DSM/DEM 정합 검증"),
        ("원우식",     "SOCAR ML/AI Engineer",           "DSM→DEM Converter 개발\nSAM3 객체 추출 · GDAL 보간 · QA"),
        ("한영재",     "한화 채권 트레이더 +\n인공지능융합대학원 컴퓨팅",
                                                       "그늘막 입지 분석 전체 시스템 통합\nMCDA · CV-A/B · Self-consistency · BayesOpt · 배낭 최적화"),
    ]
    cw = [Inches(1.2), Inches(3.0), Inches(4.45)]
    tbl = slide.shapes.add_table(rows=len(rows), cols=3,
                                    left=LEFT, top=Inches(1.45),
                                    width=sum(cw, Inches(0)),
                                    height=Inches(3.0)).table
    for i, w in enumerate(cw):
        tbl.columns[i].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
            cell.margin_left = Inches(0.1)
            tf = cell.text_frame; tf.word_wrap = True
            for li, line in enumerate(val.split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                if ri == 0:
                    _run(p.add_run(), line, size=11, bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF))
                    cell.fill.solid(); cell.fill.fore_color.rgb = C_TEXT
                else:
                    color = C_ACCENT if ci == 0 else (C_SUB if ci == 1 else C_TEXT)
                    _run(p.add_run(), line, size=9, color=color,
                          bold=(ci == 0))

    # 하단: 협업 방식
    _box(slide, LEFT, Inches(4.65), CONTENT_WIDTH, Inches(0.6),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(0.85), Inches(4.75), Inches(8.3), Inches(0.3),
         "💡 협업: 두 트랙 독립 개발 → 공통 자산(DSM·SAM3 마스크) 공유 → 최종 발표 시 통합",
         size=10, color=C_TEXT)
    _footer(slide, 3)


def build_timeline(slide):
    """[4p] 일정·마일스톤."""
    _header(slide, "일정 · 마일스톤",
            "4월 착수 → 6월 4일 중간 → 6월 18일 최종 발표")

    # 간트 차트 스타일 (간단 막대)
    phases = [
        ("4월",        "착수·역할 분담",          0.0, 0.15, C_SUB),
        ("4~5월",      "트랙 1 개발 (DSM→DEM)",  0.10, 0.55, C_ACCENT),
        ("4~5월",      "트랙 2 개발 (그늘막)",    0.10, 0.65, C_GREEN),
        ("5월 말",     "DSM 자산 공유·통합",      0.55, 0.75, C_AMBER),
        ("6/4",       "중간 점검 발표 (오늘)",    0.70, 0.78, C_RED),
        ("6/4~6/18",   "보완·문서·발표 준비",     0.78, 1.0, C_SUB),
    ]
    base_x = Inches(2.5)
    base_y = Inches(1.5)
    track_h = Inches(0.35)
    track_gap = Inches(0.50)
    bar_full = Inches(6.5)

    # 시간 라벨
    _txt(slide, Inches(2.5), Inches(1.3), Inches(0.5), Inches(0.2),
         "4월", size=8, color=C_MUTED)
    _txt(slide, Inches(5.7), Inches(1.3), Inches(0.5), Inches(0.2),
         "5월 말", size=8, color=C_MUTED)
    _txt(slide, Inches(8.7), Inches(1.3), Inches(0.5), Inches(0.2),
         "6/18", size=8, color=C_MUTED)

    for i, (period, label, start, end, color) in enumerate(phases):
        y = base_y + i * track_gap
        # 라벨
        _txt(slide, LEFT, y + Inches(0.05), Inches(1.6), Inches(0.3),
             label, size=10, color=C_TEXT)
        _txt(slide, LEFT, y + Inches(0.28), Inches(1.6), Inches(0.2),
             period, size=7, color=C_MUTED)
        # 배경 바
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, base_x, y, bar_full, track_h,
        )
        bg.fill.solid(); bg.fill.fore_color.rgb = C_BORDER
        bg.line.fill.background()
        # 진행 바
        bar_x = base_x + Emu(int(bar_full * start))
        bar_w = Emu(int(bar_full * (end - start)))
        if bar_w > 0:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, bar_x, y, bar_w, track_h,
            )
            bar.fill.solid(); bar.fill.fore_color.rgb = color
            bar.line.fill.background()

    # 오늘 위치 마커
    today_x = base_x + Emu(int(bar_full * 0.72))
    today = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      today_x, base_y - Inches(0.1),
                                      Emu(38100), Inches(3.4))
    today.fill.solid(); today.fill.fore_color.rgb = C_RED
    today.line.fill.background()

    _txt(slide, today_x - Inches(0.4), base_y - Inches(0.35),
         Inches(0.9), Inches(0.25),
         "▼ 오늘 (6/4)", size=8, bold=True, color=C_RED)

    _footer(slide, 4)


def build_track1(slide):
    """[5p] 트랙 1 진행 상황."""
    _header(slide, "트랙 1 — DSM → DEM Converter",
            "정사영상에서 객체 추출 → 보간 → 순수 지형 모델")

    # 진행도 박스
    _box(slide, LEFT, Inches(1.45), CONTENT_WIDTH, Inches(0.55),
          fill=C_GREEN, border=None)
    _txt(slide, Inches(0.85), Inches(1.58), Inches(8.3), Inches(0.3),
         "✓ 진행 상황: 5단계 파이프라인 완성 + UI/UX 완료 + 솔루션 연계 흐름 검증",
         size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # 5 STEP 카드
    steps = [
        ("STEP 1", "입력·정합",   "raster_io · 표준화 · 정합 검사"),
        ("STEP 2", "SAM3 추출",  "건물·수목·자동차 마스크"),
        ("STEP 3", "객체 제거",  "신뢰도 임계 + 마스크 융합"),
        ("STEP 4", "GDAL 보간", "탐색 100px + 평활화"),
        ("STEP 5", "QA·저장",   "정량 지표 + GeoTIFF"),
    ]
    card_w = Inches(1.68); card_h = Inches(1.65)
    for i, (tag, name, body) in enumerate(steps):
        x = LEFT + Inches(i * (1.68 + 0.07))
        c = _box(slide, x, Inches(2.2), card_w, card_h,
                  fill=C_PANEL, border=C_BORDER, border_w=0.75)
        tf = c.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.13)
        p = tf.paragraphs[0]
        _run(p.add_run(), tag, size=8, bold=True, color=C_ACCENT)
        p2 = tf.add_paragraph()
        _run(p2.add_run(), name, size=11, bold=True, color=C_TEXT)
        p2.space_before = Pt(2)
        p3 = tf.add_paragraph()
        _run(p3.add_run(), body, size=7, color=C_SUB)
        p3.space_before = Pt(3)

    # 기술 스택
    _box(slide, LEFT, Inches(4.05), Inches(4.15), Inches(1.15),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(0.85), Inches(4.15), Inches(4), Inches(0.3),
         "🛠 기술 스택", size=10, bold=True, color=C_ACCENT)
    _txt(slide, Inches(0.85), Inches(4.40), Inches(4), Inches(0.3),
         "PyQt6 · SAM3 · GDAL · rasterio · folium",
         size=9, color=C_TEXT, font=FONT_MONO)
    _txt(slide, Inches(0.85), Inches(4.65), Inches(4), Inches(0.4),
         "확장: SAM3 파인튜닝 → 추출 정확도 향상",
         size=8, color=C_SUB)

    # 다음 단계
    _box(slide, Inches(5.18), Inches(4.05), Inches(4.15), Inches(1.15),
          fill=RGBColor(0xFF, 0xF8, 0xE6), border=C_AMBER)
    _txt(slide, Inches(5.35), Inches(4.15), Inches(4), Inches(0.3),
         "📋 남은 작업", size=10, bold=True, color=C_AMBER)
    _bullets(slide, Inches(5.35), Inches(4.40), Inches(4), Inches(0.8), [
        "테스트 데이터셋 확장",
        "SAM3 파인튜닝 실험",
        "최종 발표용 데모 시나리오",
    ], size=8)
    _footer(slide, 5)


def build_track2(slide):
    """[6p] 트랙 2 진행 상황."""
    _header(slide, "트랙 2 — 그늘막 입지 분석",
            "DSM 재활용 + 다기준 의사결정으로 최적 입지 추천")

    _box(slide, LEFT, Inches(1.45), CONTENT_WIDTH, Inches(0.55),
          fill=C_GREEN, border=None)
    _txt(slide, Inches(0.85), Inches(1.58), Inches(8.3), Inches(0.3),
         "✓ 진행 상황: 7 피처 MCDA + CV-A/B + Self-consistency + BayesOpt + 배낭 최적화 통합 완료",
         size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # 좌측: 모듈 진행도
    _txt(slide, LEFT, Inches(2.15), Inches(4.2), Inches(0.3),
         "구현 완료 모듈 (12개)", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, LEFT, Inches(2.45), Inches(4.2), Inches(2.7), [
        "STEP 0 격자 생성 (3,672 cells)",
        "STEP 1 Score 7 피처 가중합 (MCDA)",
        "STEP 2 공간 필터 4단계 (보행로·건물·결집지)",
        "STEP 3 6 시나리오 민감도 → 강건 입지",
        "CV-A V-World + Mobile-SAM (건물 30동)",
        "CV-B Mapillary + SegFormer (그늘 결핍)",
        "CV-DSM 4시점 ray-cast (Deep Umbra 영감)",
        "Self-consistency 5회 (강의 40p)",
        "Bayesian Opt 가중치 자동 튜닝 (NEW)",
        "PuLP 예산 배낭 최적화",
    ], size=8)

    # 우측: 핵심 결과
    _box(slide, Inches(5.18), Inches(2.15), Inches(4.15), Inches(3.0),
          fill=C_PANEL, border=C_ACCENT, border_w=1.25)
    _txt(slide, Inches(5.35), Inches(2.25), Inches(4), Inches(0.3),
         "📊 현재까지 결과", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(5.35), Inches(2.55), Inches(4), Inches(2.5), [
        "동작구 3,672 격자 → 330 후보",
        "강건 입지 2곳 (동작대로 사당-이수)",
        "흑석동 사각지대 4곳 식별",
        "  · TOP4 최대 359.6m (기존 그늘막 대비)",
        "예산 4,000만원 → 5개 배낭 최적 배치",
        "  · status=Optimal, 200m 분산 만족",
        "학습 가중치: shade −0.15→−0.30",
        "  · 실측 18개 검증 부합 (Δ−0.21)",
    ], size=8)
    _footer(slide, 6)


def build_linkage(slide):
    """[7p] 두 트랙 연계와 학제간 협업."""
    _header(slide, "두 트랙의 연계",
            "같은 DSM 자산을 두 도구가 다른 가치로 변환")

    # 중앙: DSM (공통 자산)
    _box(slide, Inches(4.0), Inches(2.5), Inches(2.0), Inches(1.0),
          fill=C_ACCENT, border=None)
    _txt(slide, Inches(4.0), Inches(2.78), Inches(2.0), Inches(0.4),
         "DSM 원본", size=14, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _txt(slide, Inches(4.0), Inches(3.10), Inches(2.0), Inches(0.3),
         "(공통 자산)", size=9,
         color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    # 좌측: 트랙 1 산출
    _box(slide, LEFT, Inches(2.0), Inches(3.0), Inches(2.0),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(0.85), Inches(2.10), Inches(2.7), Inches(0.3),
         "트랙 1 산출", size=10, bold=True, color=C_SUB)
    _bullets(slide, Inches(0.85), Inches(2.40), Inches(2.7), Inches(1.5), [
        "SAM3 건물 마스크",
        "GDAL 보간 알고리즘",
        "DEM (지형 분석용)",
        "QA 등급 표준",
    ], size=8)

    # 우측: 트랙 2 활용
    _box(slide, Inches(6.33), Inches(2.0), Inches(3.0), Inches(2.0),
          fill=C_PANEL, border=C_ACCENT, border_w=1.25)
    _txt(slide, Inches(6.50), Inches(2.10), Inches(2.7), Inches(0.3),
         "트랙 2 활용", size=10, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(6.50), Inches(2.40), Inches(2.7), Inches(1.5), [
        "DSM → ray-cast 그림자 입력",
        "SAM3 마스크 → CV-A 검증",
        "GDAL → 결측치 보완 응용",
        "DEM 은 사용 X (그림자 X)",
    ], size=8)

    # 화살표
    a1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(3.7), Inches(2.85),
                                   Inches(0.3), Inches(0.3))
    a1.fill.solid(); a1.fill.fore_color.rgb = C_SUB
    a1.line.fill.background()
    a2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(6.0), Inches(2.85),
                                   Inches(0.3), Inches(0.3))
    a2.fill.solid(); a2.fill.fore_color.rgb = C_ACCENT
    a2.line.fill.background()

    # 하단: 핵심 메시지
    _box(slide, LEFT, Inches(4.20), CONTENT_WIDTH, Inches(0.95),
          fill=C_TEXT, border=None)
    _txt(slide, Inches(0.85), Inches(4.30), Inches(8.3), Inches(0.3),
         "🎓 강의 13p 추천 Deep Umbra (GAN 일조) 를 ray-cast 로 단순 이식",
         size=10, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    _txt(slide, Inches(0.85), Inches(4.62), Inches(8.3), Inches(0.3),
         "→ 도시공학 + ML/AI + 도메인(트레이딩 사고법) 의 학제간 협업",
         size=10, color=RGBColor(0xFF, 0xFF, 0xFF))
    _txt(slide, Inches(0.85), Inches(4.92), Inches(8.3), Inches(0.3),
         "한 사람만으로는 만들 수 없는 결과물 — DSM 자산 재활용이 핵심 연결고리",
         size=9, color=RGBColor(0xCC, 0xCC, 0xCC))
    _footer(slide, 7)


def build_next(slide):
    """[8p] 남은 작업 + 최종 발표 준비."""
    _header(slide, "남은 작업 · 최종 발표 준비",
            "6/4 중간 점검 → 6/18 최종 발표 (자율학습 1주 포함 약 2주)")

    # 트랙 1 남은 작업
    _box(slide, LEFT, Inches(1.45), Inches(4.2), Inches(2.0),
          fill=RGBColor(0xFF, 0xF8, 0xE6), border=C_AMBER, border_w=0.75)
    _txt(slide, Inches(0.85), Inches(1.55), Inches(4), Inches(0.3),
         "트랙 1 — 남은 작업", size=11, bold=True, color=C_AMBER)
    _bullets(slide, Inches(0.85), Inches(1.85), Inches(4), Inches(1.5), [
        "테스트 데이터셋 확장 · 정확도 측정",
        "SAM3 파인튜닝 실험 (정확도 향상)",
        "최종 발표용 데모 시나리오 정리",
        "QA 등급 지표 문서화",
    ], size=9)

    # 트랙 2 남은 작업
    _box(slide, Inches(5.18), Inches(1.45), Inches(4.15), Inches(2.0),
          fill=RGBColor(0xF0, 0xF8, 0xFF), border=C_ACCENT, border_w=0.75)
    _txt(slide, Inches(5.35), Inches(1.55), Inches(4), Inches(0.3),
         "트랙 2 — 남은 작업", size=11, bold=True, color=C_ACCENT)
    _bullets(slide, Inches(5.35), Inches(1.85), Inches(4), Inches(1.5), [
        "서울 열린데이터 API 연결 (실데이터 5종)",
        "CV-D 멀티모달 VLM 시도 (강의 40~42p)",
        "동작구 전 동 실측 그늘막 수집",
        "발표 슬라이드 추가 자료 정리",
    ], size=9)

    # 통합 마일스톤
    _box(slide, LEFT, Inches(3.60), CONTENT_WIDTH, Inches(1.55),
          fill=C_PANEL, border=C_BORDER)
    _txt(slide, Inches(0.85), Inches(3.70), Inches(8.3), Inches(0.3),
         "🎯 6/18 최종 발표 통합 마일스톤", size=11, bold=True, color=C_TEXT)
    _bullets(slide, Inches(0.85), Inches(4.00), Inches(8.3), Inches(1.1), [
        "두 트랙 산출물을 한 PPTX 로 통합 (현재 26장 → 추가 정리)",
        "발표 시연: DSM→DEM 변환 (라이브) + 그늘막 추천 지도 (인터랙티브 HTML)",
        "스피커 노트 + 시연 시나리오 + 예상 질문 대비 완료",
        "조원 1인 발표 / 2인 시연 보조 — 역할 분담 확정",
    ], size=9)
    _footer(slide, 8)


# ─────────── 스피커 노트 ───────────
NOTES = {
    1: """[1p · 표지 · 60초]
- 안녕하세요. 중간 발표 시작하겠습니다.
- 데이터기반 도시설계 기말 프로젝트의 진행 상황을 공유드립니다.
- 3인 학제간 팀입니다: GIS 엔지니어 문치국, ML/AI 엔지니어 원우식, 그리고 발표자 한영재.
- 두 개의 트랙(DSM→DEM Converter + 그늘막 입지 분석)을 동시 진행 중이며,
  오늘은 전체 결과보다 진행 과정·역할 분담·향후 계획을 중심으로 말씀드리겠습니다.
""",
    2: """[2p · 개요 · 60초]
- 본 프로젝트는 두 개의 트랙으로 진행됩니다.
- 트랙 1: 친구 두 분이 만든 DSM → DEM Converter. 정사영상에서 SAM3 로 객체를 추출하고
  GDAL 로 보간해 순수 지형 모델(DEM)을 생성하는 도구입니다.
- 트랙 2: 제가 담당하는 그늘막 입지 분석. 트랙 1 에서 사용한 DSM 원본을 그대로 재활용해
  동작구 여름 그늘막 최적 입지를 추천합니다.
- 핵심 연결고리: 같은 DSM 자산을 두 도구가 다른 가치로 변환 — 학제간 협업의 핵심 메시지.
""",
    3: """[3p · 팀 구성 · 90초]
- 각자의 전문성을 모듈로 분리해 협업했습니다.
- 문치국: GIS 측량 BIM 엔지니어. NGII 흑석동 DSM·DEM·건물 shp 데이터를 확보하고
  측량 도메인 자문을 담당. 데이터 정합성을 검증해 줍니다.
- 원우식: SOCAR ML/AI Engineer. DSM→DEM Converter 의 실제 개발을 담당.
  SAM3 객체 추출, GDAL 보간, QA 지표 모두 그가 구현했습니다.
- 한영재: 금융 데이터에서 다변수 최적화를 다루던 경험을 도시 공간에 이식.
  그늘막 입지 분석의 전체 시스템 (MCDA + CV + GIS + 자동 튜닝 + 배낭 최적화) 통합.
- 협업 방식: 두 트랙 독립 개발 → 공통 자산(DSM·SAM3 마스크) 공유 → 최종 발표 시 통합.
""",
    4: """[4p · 일정 · 60초]
- 4월에 착수 후 역할을 분담했습니다.
- 4~5월 동안 두 트랙 병렬 개발. 트랙 1은 SAM3 통합과 GDAL 보간까지,
  트랙 2는 MCDA Score 식과 CV-A/B 모듈까지 완성.
- 5월 말부터 DSM 자산 공유와 통합 작업을 시작.
- 오늘 6월 4일이 중간 점검 발표.
- 앞으로 약 2주: 6/11 자율학습 주간 동안 보완 작업, 6/18 최종 발표.
- 자율학습 기간에 트랙 1은 SAM3 파인튜닝, 트랙 2는 CV-D 멀티모달 VLM 도입 계획.
""",
    5: """[5p · 트랙 1 · 75초]
- DSM→DEM Converter 는 PyQt6 데스크톱 앱으로 완성됐습니다.
- 5단계 파이프라인 모두 구현 완료:
  STEP 1 입력·정합 (raster_io 모듈)
  STEP 2 SAM3 추출 (건물·수목·자동차 마스크 + 신뢰도)
  STEP 3 객체 제거 (마스크 융합)
  STEP 4 GDAL 보간 (탐색 100px + 평활화 0~20 반복)
  STEP 5 QA 평가 + GeoTIFF 저장
- 기술 스택: PyQt6, SAM3, GDAL, rasterio, folium.
- 인터랙티브 지도 (Leafmap + Folium) 와 노션 스타일 UI 까지 완료.
- 남은 작업: 테스트 데이터셋 확장, SAM3 파인튜닝 실험, 최종 데모 시나리오 정리.
""",
    6: """[6p · 트랙 2 · 75초]
- 그늘막 입지 분석은 12개 모듈로 구성된 7 피처 MCDA 시스템입니다.
- 구현 완료: 격자 3,672 → 4단계 필터 → 330 후보 → 6 시나리오 비교 → 강건 입지.
- CV-A 는 Mobile-SAM 으로 항공사진에서 건물 30동 자동 추출.
- CV-B 는 SegFormer 로 거리뷰 1,302장 분석.
- CV-DSM 은 강의 13p 추천 Deep Umbra 영감으로 4시점 ray-cast.
- Self-consistency 5회 평균 (강의 40p) 으로 결과 안정성 검증.
- Bayesian Optimization 으로 가중치 자동 튜닝.
- PuLP 정수 선형계획법으로 예산 4천만원 배낭 최적화 — 5개 격자 선정.
- 핵심 결과: 강건 입지 2곳, 흑석동 사각지대 4곳 식별, 학습된 shade 페널티 가설 부합.
""",
    7: """[7p · 연계 · 75초]
- 두 트랙의 핵심 연결고리는 DSM 원본 자산입니다.
- 트랙 1 에서 만든 자산: SAM3 건물 마스크 (객체 추출), GDAL 보간 알고리즘, DEM (지형 분석용), QA 등급 표준.
- 트랙 2 가 활용: DSM 그대로 ray-cast 그림자 입력으로 사용,
  SAM3 마스크는 CV-A SAM 추출의 검증 자료, GDAL 은 결측치 보완에 응용.
- 단, DEM 은 트랙 2 에서 사용 안 함 — 객체 제거된 지표만 있어 그림자 시뮬레이션이 불가능.
- 교수님이 강의 13p 에 우리 프로젝트 옆에 직접 인용한 Deep Umbra 논문 (GAN 일조 컴퓨테이션) 을
  ray-cast 로 단순 이식한 것이 핵심 학술 정합성 포인트입니다.
""",
    8: """[8p · 다음 단계 · 60초]
- 6/4 중간 발표 이후 약 2주 동안 마무리 작업이 남았습니다.
- 트랙 1: 테스트 확장, SAM3 파인튜닝, 데모 시나리오 정리.
- 트랙 2: 서울 열린데이터 API 연결, CV-D 멀티모달 VLM 시도, 발표 슬라이드 정리.
- 통합 마일스톤:
  · 두 트랙 산출물을 한 PPTX 로 통합 (현재 26장 → 정리)
  · 발표 시연: DSM→DEM 라이브 + 그늘막 인터랙티브 지도 시연
  · 스피커 노트와 예상 질문 대비 완료
  · 발표 역할 분담 확정
- 6/18 최종 발표에서 더 완성된 결과로 다시 뵙겠습니다. 감사합니다.
""",
}


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    builders = [
        build_cover,
        build_overview,
        build_team,
        build_timeline,
        build_track1,
        build_track2,
        build_linkage,
        build_next,
    ]
    for builder in builders:
        slide = _new_slide(prs)
        builder(slide)

    # 노트 삽입
    for i, slide in enumerate(prs.slides, 1):
        note = NOTES.get(i, "")
        if note:
            slide.notes_slide.notes_text_frame.text = note.strip()

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"[완료] {OUT}")
    print(f"  총 {len(prs.slides)}장 · 16:9 · Noto Sans KR")
    print(f"  스피커 노트 {len(NOTES)}장 자동 삽입")


if __name__ == "__main__":
    build()
