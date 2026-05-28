"""최종 발표자료 (PPTX) 생성 스크립트 — v2 (CV-DSM + OSMnx + 흑석동 실측 통합).

실행:
    python presentation/generate_pptx.py

산출물:
    presentation/동작구_그늘막_최종발표.pptx (15장)
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "presentation" / "동작구_그늘막_최종발표.pptx"

# 색상
COLOR_PRIMARY = RGBColor(0x1E, 0x88, 0xE5)
COLOR_ACCENT  = RGBColor(0xE5, 0x39, 0x35)
COLOR_WARN    = RGBColor(0xFB, 0x8C, 0x00)
COLOR_VULN    = RGBColor(0x8E, 0x24, 0xAA)
COLOR_GREEN   = RGBColor(0x43, 0xA0, 0x47)
COLOR_CYAN    = RGBColor(0x00, 0xAC, 0xC1)
COLOR_DARK    = RGBColor(0x21, 0x21, 0x21)
COLOR_SUB     = RGBColor(0x55, 0x55, 0x55)
COLOR_LIGHT   = RGBColor(0xEE, 0xEE, 0xEE)
FONT = "맑은 고딕"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 16


# 슬라이드별 스피커 노트 (PPTX notes_slide 에 자동 삽입)
# 상세 버전은 presentation/speaker_notes.md 참고
SLIDE_NOTES = {
    1: """[표지 · 30초]
- 인공지능융합대학원 한영재 / 2025961227
- 후크: 그늘막 입지를 민원·직관이 아닌 데이터의 논리로 결정
- MCDA 코어 + CV 3종 + GIS 2종 = 9개 외부 데이터 통합
- 코드: github.com/thgtot92/project_1
""",
    2: """[문제 정의 · 접근 · 45초]
- 핵심 메시지: 데이터 통합 + 다기준 의사결정으로 객관적 최적 입지 도출
- 기존 방식 한계: 민원 기반 → 이미 커버된 곳 중복, 사각지대 누락
- 우리 방식: 9종 데이터 → 7 피처 가중합 → 다단계 필터 → 6 시나리오 비교 → 강건 입지
- 중요: CV·GIS는 Score를 대체하지 않음. 입력 데이터 품질만 끌어올림. 설명가능성 유지

예상 질문: 왜 딥러닝으로 Score 학습 안 했나? → 정답 라벨 없음 + 설명가능성 의도적 선택
""",
    3: """[파이프라인 개요 · 60초]
- 5단계: 격자(3672) → Score(7피처 가중합) → 필터(330) → 시나리오(6) → 근거+시각화
- 첫 실행 ~10분, 캐시 이후 ~30초
- 9종 데이터, 7 피처, 4단계 필터, 6 시나리오 — 핵심 수치 강조
""",
    4: """[입력 데이터 9종 · 60초]
- 3 계층 구조 강조: 기본 5종 (CSV·shp), CV 2종 (외부 API), 외부 정밀 2종 (NGII·OSM)
- 5종 기본 데이터 중 4종 더미. data_loader.py 가 실데이터 자동 전환
- 외부 API: 첫 1회 호출 후 디스크 캐시. 발표 도중 호출 X

예상 질문: 더미로 발표 신뢰성? → 동작구 실제 거점 기반 가우시안. 실데이터 자동 전환 구조라 인프라 검증용 충분
""",
    5: """[Score 식 7 피처 · 60초]
- 한 줄 공식: Score = 0.18 pop + 0.18 lst + 0.18 vuln − 0.15 shade − 0.05 natural + 0.12 sv + 0.20 inter
- 양수 합 0.86, 음수 -0.20. MinMax [0,1] 정규화 후 가중합
- 시나리오마다 가중치 재배분. 가중치 절대정답은 없음
- Streamlit 슬라이더로 즉시 조정 가능 (다음 다음 슬라이드에서 시연)

예상 질문: 가중치 어떻게 정함? → 도메인 직관 + 시나리오 민감도. 6개 시나리오 모두에서 살아남는 입지가 강건 입지
""",
    6: """[공간 필터링 4단계 · 60초]
- 단순 Score만으론 건물 위·녹지에 점 같은 시각적 오류 발생
- 4단계 컷: 보행로 5m → 그늘막 150m 외 → 건물 안 컷 → 결집지 50m 근접
- 3,672 → 449 → 446 → 434 → 330 (약 9%만 살아남음)
- 슬라이드 6 의 지도(shade_map.html)는 실시간 임베드. 필요시 직접 클릭/줌 시연
""",
    7: """[CV-A · SAM · 60초]
- V-World WMTS z=15 → 48 타일 합성 (2048×1536)
- Mobile-SAM (Meta 2023): 40MB CPU 추론, zero-shot
- 필터: 면적·종횡비·밝기 3중 → 30동 추출
- 흑석동에서는 NGII 3,775동이 자동 union → SAM 30동의 125배 정밀

예상 질문: YOLO 대신 SAM 이유? → zero-shot, 라벨 불필요. YOLO 는 학습 데이터 비용 큼
""",
    8: """[CV-B · SegFormer · 60초]
- Mapillary BBOX 4×4 검색 → 1,302장 발견. Meta 소유 오픈 거리뷰
- 후보 19격자에 nearest 매핑, 격자당 ≤3장
- SegFormer-b0 (HuggingFace) CityScapes 19 classes pretrained
- 그늘 결핍 = (보도+도로) × (1 − 건물 − 식생)
- TOP10 평균 deficit 0.153 → Score 6번째 피처
- 선행연구: MIT Treepedia, Place Pulse 2.0

예상 질문: 거리뷰 시기 다른데? → 정적 요소 비율은 안정적
""",
    9: """[CV-DSM · Deep Umbra 영감 · 75초 — 강의 13p 핵심]
- 강의자료 13p 에 우리 프로젝트 옆에 Deep Umbra 논문 박혀 있었음 (evl.uic.edu/shadows)
- Deep Umbra: GAN 으로 일조/그림자 누적 맵 생성. RMSE 0.063, SSIM 0.90
- 우리는 풀 GAN 학습 대신 DSM 실측 + 4시점 ray-cast 로 단순화
- DSM(표면) − DEM(지표) = nDSM(객체 높이)
- 4시점 (10·12·14·16시) shadow union → 누적 비율 [0,1]
- 흑석동 309격자에 평균 누적 0.157 → 단일 시점보다 훨씬 정밀
- NGII 흑석동 실측 건물 3,775동도 동시 통합

예상 질문: 왜 풀 GAN 안 썼나? → 학습 시간/데이터 부담. 4시점 ray-cast 로 같은 메시지. GAN 으로 확장 여지
""",
    10: """[OSMnx · 교차로 + 횡단보도 · 60초]
- OSMnx: OpenStreetMap Overpass API Python 래퍼
- walkable highway 그래프 → street_count ≥ 3 노드 = 교차로
- 한국 OSM 의 함정: 횡단보도가 highway=crossing 보다 footway=crossing 우세
- 처음에 흑석동 횡단보도 0개로 떴다가, 다중 태그 검색으로 해결
- highway · footway · crossing 3종 통합 → 306개 횡단보도 발견
- 결집지점 union → 격자별 80m 내 개수(피처) + 50m 근접(필터)
- 효과: 노량진역-수산시장이 처음으로 TOP10 진입

예상 질문: OSM 신뢰성? → 인구 밀집 지역은 충실. NGII 도로대장 비교는 향후 과제
""",
    11: """[시나리오 6 + 강건 입지 · 75초]
- 6 시나리오: 기본/고령자/폭염/유동인구/보행환경/교차로
- 동일 피처 위에 가중치만 바꿔 재스코어 = 정책 민감도 분석
- 유동인구 시나리오 최고 0.648 score
- 강건 입지 2곳: 동작대로 사당-이수 축 (37.4907, 126.9647) + 사당역 인근 (37.4898, 126.9670)
- "정책 관점이 바뀌어도 정답이 안 바뀌는 곳" = 예산 제약 최우선
- 흑석동 한강변은 이전 강건 입지였으나 교차로 시나리오 추가로 탈락 → 더 엄격한 검증 통과한 사당-이수의 강건성 ↑

예상 질문: 2곳만이면 너무 적지 않나? → 6 시나리오 합의 자체가 어려운 일. 2곳 살아남았다는 게 의미
""",
    12: """[1차 발표 3대 구역 검증 · 45초]
- 1차 발표 때 노량진·사당-이수·상도로 3대 후보 제시
- 알고리즘이 얼마나 재현했는지 자기 검증
- 노량진: △ → ✓ (횡단보도 통합 후 신규 진입)
- 사당-이수: ✓ TOP1·2·4·5 차지
- 상도로: ✓ → △ (보행로 5m 강화로 컷, 정책적 재검토 필요)
""",
    13: """[흑석동 실측 vs 신규 추천 · 75초 — 핵심 임팩트]
- 이 슬라이드가 이번 프로젝트 최대 임팩트 포인트
- 사용자 제공 실측 그늘막 18개 (고정 13 + 스마트 5, 2019~2023)와 알고리즘 TOP10 비교
- 3 카테고리 분류: 사각지대(>150m) / 중간 / 보강(<60m)
- 사각지대 4곳: TOP3·4·5·8
- TOP4 (37.49786, 126.96129) 최대 사각지대 359.6m — 정책 어필 1순위
- 정책: 흑석동은 그늘막 75% 영역 커버, 다른 동 우선. 흑석동 내부 추가는 사각지대 4곳

예상 질문: 사각지대인데 score 낮으면 추천 가치? → 그늘막 0이라 shade 페널티 0. 거리뷰·DSM 보강 시 score 같이 상승 가능
""",
    14: """[예산 제약 최적화 · 60초 — NEW]
- 정책 직답: "예산 N원 = 그늘막 K개 = 어디?"
- PuLP 정수 선형계획법 (CBC 솔버, 무료)
- 목표: max Σ score_i · x_i
- 제약: 예산(Σ cost·x ≤ B) + 공간 분산(가까운 격자 동시선택 X)
- 후보 풀 50개 → 200m 이격 제약 29쌍 active
- 결과: status=Optimal, 5개 선정, 총 score 1.852, 예산 100% 사용
- 선정 분포: 사당-이수 4 + 노량진 1 (자동 분산)
- TOP10 단순 정렬과 다른 점: 200m 이격으로 같은 곳 몰림 방지

예상 질문: 단가 800만원 근거? → 더미. 실 단가 확보 시 즉시 교체. 단가가 달라도 비례적으로 K개 결정
""",
    15: """[Streamlit 대시보드 · 45초 — 시연용]
- streamlit run app.py → http://localhost:8501
- 사이드바: 시나리오 프리셋 + 7개 가중치 슬라이더
- 본문: Folium 지도 + TOP10 표 + 통계
- 캐싱: 격자·피처 1회, 슬라이더는 가중합만 (밀리초)
- 발표 중 "vuln 더 올리면?" 질문에 즉답
- 시간 여유 있으면 별도 창에서 슬라이더 1~2개 움직여 시연
""",
    16: """[한계 + 한 줄 요약 · 45초]
- 한계: 4종 기본 데이터 더미 / LST 래스터 미연결 / SAM 추정 높이 / 단가 800만원 가정
- 다음 단계: API 연결 / 전 동 실측 / Streamlit 시각화 / 타 자치구 확장 (BBOX 교체)
- 한 줄 요약: 항공·거리뷰·DSM·도로망·실측 + MCDA + 배낭 최적화로 강건 2 + 사각지대 4 + 예산 5곳
- "감사합니다. 질문 받겠습니다."

예상 질문:
- 타 자치구 적용? → BBOX·NGII 교체. 가중치 재튜닝 필요
- 코드? → github.com/thgtot92/project_1 (MIT)
""",
}


def _set_font(run, size=18, bold=False, color=COLOR_DARK, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(slide, left, top, width, height, text,
              size=18, bold=False, color=COLOR_DARK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    return tx


def _add_bullets(slide, left, top, width, height, items, size=16, bullet="• "):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = bullet + item
        _set_font(run, size=size, color=COLOR_DARK)
        p.space_after = Pt(5)
    return tx


def _add_header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5),
        Inches(0.12), Inches(0.8),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    _add_text(slide, Inches(0.75), Inches(0.45),
              Inches(12), Inches(0.6),
              title, size=28, bold=True, color=COLOR_DARK)
    if subtitle:
        _add_text(slide, Inches(0.75), Inches(1.0),
                  Inches(12), Inches(0.4),
                  subtitle, size=13, color=COLOR_SUB)


def _add_footer(slide, page_num):
    _add_text(slide, Inches(0.5), Inches(7.0),
              Inches(10), Inches(0.4),
              "동작구 여름 그늘막 최적 입지 추천 시스템 · 한영재 · 2025961227",
              size=10, color=COLOR_SUB)
    _add_text(slide, Inches(12.3), Inches(7.0),
              Inches(1), Inches(0.4),
              f"{page_num} / {TOTAL}", size=10,
              color=COLOR_SUB, align=PP_ALIGN.RIGHT)


def _box(slide, left, top, width, height, fill, border=None, border_w=1.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if border is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = border
        box.line.width = Pt(border_w)
    return box


def _placeholder_image(slide, left, top, width, height, caption):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT
    box.line.color.rgb = COLOR_SUB
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"📷 {caption}"
    _set_font(run, size=13, bold=True, color=COLOR_SUB)


def _table(slide, left, top, total_width, height, rows, col_widths,
            header_color=COLOR_PRIMARY, header_size=13, body_size=12,
            bold_rows=None):
    tbl = slide.shapes.add_table(rows=len(rows), cols=len(rows[0]),
                                   left=left, top=top,
                                   width=total_width, height=height).table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw
    bold_rows = bold_rows or []
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            if ri == 0:
                _set_font(run, size=header_size, bold=True,
                           color=RGBColor(0xFF, 0xFF, 0xFF))
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                _set_font(run, size=body_size,
                           bold=(ri in bold_rows), color=COLOR_DARK)
    return tbl


# ───────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ────── 1. 표지 ──────
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_DARK
    bg.line.fill.background()

    _add_text(s, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.0),
              "동작구 여름 그늘막 최적 입지 추천 시스템",
              size=38, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(s, Inches(1.0), Inches(2.9), Inches(11.3), Inches(0.6),
              "MCDA × 컴퓨터비전(SAM·SegFormer) × DSM(Deep Umbra) × OSMnx",
              size=17, color=RGBColor(0xBB, 0xBB, 0xBB))
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1.0), Inches(4.0),
                               Inches(2.0), Emu(50800))
    line.fill.solid(); line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()

    _add_text(s, Inches(1.0), Inches(4.3), Inches(11), Inches(0.4),
              "데이터기반 도시설계 · 기말 프로젝트",
              size=16, color=RGBColor(0xCC, 0xCC, 0xCC))
    _add_text(s, Inches(1.0), Inches(4.8), Inches(11), Inches(0.4),
              "한영재 / 인공지능융합대학원 인공지능컴퓨팅 / 2025961227",
              size=14, color=RGBColor(0xAA, 0xAA, 0xAA))
    _add_text(s, Inches(1.0), Inches(5.4), Inches(11), Inches(0.4),
              "Repo: github.com/thgtot92/project_1",
              size=12, color=RGBColor(0x88, 0xAE, 0xE5))
    _add_text(s, Inches(1.0), Inches(5.9), Inches(11), Inches(0.4),
              "2026년 5월",
              size=12, color=RGBColor(0x88, 0x88, 0x88))

    # ────── 2. 문제 정의·접근 (MCDA 프레임) ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "문제 정의 · 접근",
                    "어디에 설치할 것인가 — 민원·직관이 아닌 데이터의 논리로")
    _add_bullets(s, Inches(0.75), Inches(1.7), Inches(12), Inches(2.5), [
        "그늘막은 유한한 예산 → 어느 격자에 '폭염스트레스 × 유동인구 × 취약계층' 효용 최대인가?",
        "기존 방식: 민원·직관 → 공간적 편향, 이미 커버된 곳 중복 설치",
        "제안: 다변수 가중합 + 다단계 공간 필터 + 다정책 시나리오 + 4종 외부 데이터 통합",
    ], size=15)

    box = _box(s, Inches(0.75), Inches(4.2), Inches(12), Inches(2.3),
                RGBColor(0xF5, 0xF8, 0xFD), COLOR_PRIMARY)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "다기준 의사결정 (MCDA) 기반 입지 추천"
    _set_font(r, size=19, bold=True, color=COLOR_PRIMARY)
    for t in [
        "• 다변수 가중합 + Min-Max 정규화로 객관적 점수 산출",
        "• 공간 제약 4단계 (보행로·건물·교차로·기존그늘막)로 실제 설치 가능 위치만",
        "• 6 시나리오 민감도 분석 → 정책 관점 불변 강건 입지 식별",
        "• 컴퓨터비전·GIS는 'Score를 대체'하는 게 아니라 '입력 데이터 품질을 끌어올림'",
    ]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = t
        _set_font(r, size=14, color=COLOR_DARK)
        p.space_before = Pt(3)
    _add_footer(s, 2)

    # ────── 3. 파이프라인 개요 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "파이프라인 개요",
                    "9종 데이터 → 7 피처 → 4단계 필터 → 6 시나리오 → TOP 10")

    steps = [
        ("STEP 0", "격자", "동작구 100m 격자\n3,672 cells", COLOR_SUB),
        ("STEP 1", "Score", "7 피처 가중합\nMinMax 정규화", COLOR_PRIMARY),
        ("STEP 2", "필터", "보행로 5m → 그늘막 →\n건물컷 → 결집지 50m\n3,672 → 330", COLOR_WARN),
        ("STEP 3", "시나리오", "6 프리셋 × TOP10\n강건입지 식별", COLOR_ACCENT),
        ("STEP 4", "근거+시각화", "LLM 자연어 근거 +\nFolium 지도", COLOR_GREEN),
    ]
    left0 = Inches(0.4); top0 = Inches(2.0)
    w = Inches(2.4); h = Inches(3.3); gap = Inches(0.1)
    for i, (tag, title, body, color) in enumerate(steps):
        x = left0 + Inches(i * (2.4 + 0.1))
        box = _box(s, x, top0, w, h, RGBColor(0xFA, 0xFA, 0xFA), color, 2)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = tag
        _set_font(r, size=11, bold=True, color=color)
        p2 = tf.add_paragraph()
        r = p2.add_run(); r.text = title
        _set_font(r, size=20, bold=True, color=COLOR_DARK)
        p2.space_before = Pt(3)
        for line_txt in body.split("\n"):
            p = tf.add_paragraph()
            r = p.add_run(); r.text = line_txt
            _set_font(r, size=12, color=COLOR_SUB)
            p.space_before = Pt(2)
        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + w - Inches(0.02), top0 + Inches(1.45),
                                      Inches(0.18), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = COLOR_SUB
            arr.line.fill.background()

    _add_text(s, Inches(0.75), Inches(5.7), Inches(12), Inches(0.5),
              "외부 데이터: V-World 위성 · Mapillary 거리뷰 · 흑석동 DSM/DEM/NGII · OSMnx 도로망/교차로/횡단보도 · 흑석동 실측 그늘막 18개",
              size=12, color=COLOR_SUB)
    _add_footer(s, 3)

    # ────── 4. 데이터 9종 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "입력 데이터 (9종, 3 계층)",
                    "기본 5종 + 컴퓨터비전 2종 + 흑석동 정밀·OSMnx 2종")
    rows = [
        ("구분", "데이터", "출처", "용도"),
        ("기본", "생활인구 (시간대별)", "SKT 서울열린데이터", "popdens 피처"),
        ("기본", "지표면 온도 LST", "서울연구원 Landsat", "lst 피처"),
        ("기본", "취약계층 비율", "KOSIS 동별", "vuln 피처"),
        ("기본", "기존 그늘막 위치", "서울열린데이터 / 사용자제공", "shade 페널티 + 비교"),
        ("기본", "인도·횡단보도", "국가공간정보포털 (현재 더미)", "보행로 입력"),
        ("CV-A", "V-World 항공사진 WMTS", "V-World API", "Mobile-SAM → 건물 30동"),
        ("CV-B", "Mapillary 거리뷰", "Mapillary API", "SegFormer → 그늘 결핍"),
        ("흑석동", "NGII DSM·DEM·건물 3,775동·경계", "국토지리정보원", "DSM 누적 그림자·실측 건물"),
        ("OSMnx", "OSM 도로망·교차로·횡단보도", "Overpass API (OSMnx)", "보행로 필터·결집지 밀도"),
    ]
    widths = [Inches(1.2), Inches(3.5), Inches(3.8), Inches(3.5)]
    _table(s, Inches(0.4), Inches(1.8), sum(widths, Inches(0)),
            Inches(4.5), rows, widths, header_size=12, body_size=11)
    _add_text(s, Inches(0.75), Inches(6.4), Inches(12), Inches(0.4),
              "⚙ 5종 기본 데이터는 더미↔실데이터 자동 스위칭 · 외부 API는 첫 1회 호출 후 디스크 캐시",
              size=11, color=COLOR_SUB)
    _add_footer(s, 4)

    # ────── 5. Score 식 7 피처 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "알고리즘 · Score 식 (7 피처)",
                    "각 피처 MinMax [0,1] 정규화 후 가중합")
    box = _box(s, Inches(0.75), Inches(1.8), Inches(12), Inches(1.6),
                COLOR_DARK)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Score = 0.18·popdens + 0.18·lst + 0.18·vuln"
    _set_font(r, size=18, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = "− 0.15·shade − 0.05·natural + 0.12·streetview_deficit + 0.20·intersection_density"
    _set_font(r, size=15, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))

    _add_bullets(s, Inches(0.75), Inches(3.6), Inches(12), Inches(2.5), [
        "popdens (+0.18) — 시간대 가중 유동인구 (9~18시, 오후 1~3시 1.0 피크)",
        "lst (+0.18) — 지표면 온도, 오후 피크 가중",
        "vuln (+0.18) — 고령자·어린이 비율",
        "shade (−0.15) — 기존 그늘막 반경 150m 내 커버리지 (이미 된 곳 감점)",
        "natural (−0.05) — CV-A SAM 30동 + NGII 흑석동 3,775동 + 흑석동 DSM 4시점 누적 (Deep Umbra)",
        "streetview_deficit (+0.12) — CV-B SegFormer 거리뷰: (보도+도로) × (1−건물−식생)",
        "intersection_density (+0.20) — OSMnx 교차로 + 횡단보도 반경 80m 내 개수",
    ], size=13)

    _add_text(s, Inches(0.75), Inches(6.4), Inches(12), Inches(0.4),
              "🎚 Streamlit 대시보드(app.py)에서 슬라이더로 실시간 조정 + 시나리오 6개 프리셋 내장",
              size=12, bold=True, color=COLOR_PRIMARY)
    _add_footer(s, 5)

    # ────── 6. 공간 필터링 강화 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "공간 필터링 — 4단계 다중 컷",
                    "잘못된 위치(차도·건물·녹지·외진곳)에 그늘막 추천 방지")

    funnel = [
        ("3,672", "전체 격자", COLOR_SUB),
        ("449", "보행로 5m 이내\n(OSMnx)", COLOR_WARN),
        ("434", "건물 안 컷\n(NGII 5m inset)", COLOR_ACCENT),
        ("330", "교차로·횡단보도\n50m 근접", COLOR_GREEN),
        ("10", "최종 추천", COLOR_PRIMARY),
    ]
    left = Inches(0.4); top = Inches(2.0)
    w = Inches(2.4); h = Inches(2.5)
    for i, (num, label, color) in enumerate(funnel):
        x = left + Inches(i * (2.4 + 0.1))
        box = _box(s, x, top, w, h, color)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        _set_font(r, size=38, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run(); r.text = label
        _set_font(r, size=11, color=RGBColor(0xFF, 0xFF, 0xFF))

    _add_bullets(s, Inches(0.75), Inches(4.8), Inches(12), Inches(1.8), [
        "이전 버전: 보행로 20m + 더미 14개 라인 → '건물 위'·'녹지에 점' 시각적 오류",
        "수정: OSMnx walkable edges 통합 + buffer 5m + NGII 건물 안 centroid 컷 + 결집지 50m 근접",
        "효과: 보행 가능 + 건물 밖 + 교차로/횡단보도 옆에만 추천 (실용성 ↑)",
    ], size=13)

    _placeholder_image(s, Inches(0.75), Inches(6.4), Inches(12), Inches(0.5),
                       "output/shade_map.html 또는 scenarios_map.html 캡처 가능")
    _add_footer(s, 6)

    # ────── 7. CV-A: V-World + Mobile-SAM ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "CV-A · 항공사진 + Mobile-SAM",
                    "딥러닝 zero-shot segmentation으로 동작구 건물 자동 추출")
    _add_bullets(s, Inches(0.75), Inches(1.8), Inches(6.5), Inches(4.5), [
        "V-World WMTS z=15 → 동작구 BBOX 48 타일 합성 (2048×1536)",
        "Mobile-SAM (Meta, 2023) — 40MB 경량 SAM, CPU 추론 가능",
        "필터: 면적 200~30k px², 종횡비 0.2~5, 평균 밝기 80~220",
        "결과: 실측 30동 자동 추출 → buildings.geojson",
        "픽셀 → EPSG:3857 → WGS84 polygon 변환 + 면적 기반 추정 높이",
        "흑석동에서는 NGII 실측 3,775동이 SAM 30동을 union (data_loader 자동)",
    ], size=13)
    box = _box(s, Inches(7.5), Inches(1.8), Inches(5.3), Inches(4.5),
                RGBColor(0xF5, 0xF8, 0xFD), COLOR_PRIMARY)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "📁 output/cv_buildings_overlay.png"
    _set_font(r, size=13, bold=True, color=COLOR_PRIMARY)
    for t in ["", "2048×1536 위성 이미지", "+ SAM 마스크 노란 오버레이",
              "+ 추출 polygon 녹색 외곽선", "",
              "30동 분포: 노량진·사당·이수·",
              "상도로·장승배기·흑석동"]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = t
        _set_font(r, size=11, color=COLOR_DARK)
        p.space_before = Pt(2)
    _add_footer(s, 7)

    # ────── 8. CV-B: Mapillary + SegFormer ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "CV-B · 거리뷰 + SegFormer (CityScapes)",
                    "지면 시점에서 본 보행 환경을 19-class 의미 분할")
    _add_bullets(s, Inches(0.75), Inches(1.8), Inches(6.5), Inches(4.5), [
        "Mapillary — 동작구 BBOX 4×4 분할 검색 → 1,302장 발견",
        "후보 19개 격자에 nearest 매핑 → 격자당 ≤3장 다운로드",
        "HuggingFace SegFormer-b0 (CityScapes 19 classes pretrained)",
        "핵심 클래스: building · vegetation · road · sidewalk · sky",
        "선행연구: MIT Treepedia, Place Pulse 2.0 응용",
    ], size=13)
    box = _box(s, Inches(7.5), Inches(1.8), Inches(5.3), Inches(2.0),
                COLOR_DARK)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "그늘 결핍 지수"
    _set_font(r, size=14, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = "(road + sidewalk)"
    _set_font(r, size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run(); r.text = "× (1 − building − vegetation)"
    _set_font(r, size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    box2 = _box(s, Inches(7.5), Inches(4.0), Inches(5.3), Inches(2.3),
                 RGBColor(0xFF, 0xF4, 0xE5), COLOR_WARN)
    tf = box2.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "📊 결과"
    _set_font(r, size=13, bold=True, color=COLOR_WARN)
    for t in ["TOP10 평균 deficit: 0.153",
              "Score 6번째 피처 통합",
              "신규 시나리오: 보행환경_중시"]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = "• " + t
        _set_font(r, size=12, color=COLOR_DARK)
        p.space_before = Pt(3)
    _add_footer(s, 8)

    # ────── 9. CV-DSM (Deep Umbra) + 흑석동 NGII ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "CV-DSM · 흑석동 정밀 그림자 (강의자료 13p 추천)",
                    "Deep Umbra 영감 — DSM 4시점 ray-cast 누적 그림자")
    _add_bullets(s, Inches(0.75), Inches(1.8), Inches(7.0), Inches(4.5), [
        "강의자료 13p 추천 논문: Deep Umbra (evl.uic.edu/shadows)",
        "  → GAN 기반 일조/그림자 접근성 컴퓨테이션",
        "본 프로젝트: 풀 GAN 대신 DSM 실측 + 4시점 ray-cast 단순화",
        "DSM(표면) − DEM(지표) = nDSM(객체 높이 모델)",
        "4시점 (오전 10·정오·오후 2·오후 4) shadow union → 누적 비율",
        "흑석동 309격자에 적용 → 평균 누적 그림자 0.157",
        "NGII 흑석동 실측 건물 3,775동 (층수·HEGT) 통합 (데이터 품질 ↑)",
    ], size=13)
    box = _box(s, Inches(8.0), Inches(1.8), Inches(4.8), Inches(4.5),
                RGBColor(0xF5, 0xF8, 0xFD), COLOR_PRIMARY)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "📐 처리 흐름"
    _set_font(r, size=13, bold=True, color=COLOR_PRIMARY)
    for t in ["",
              "1. rasterio 로 .tif 로딩",
              "2. nDSM = DSM − DEM (>0)",
              "3. 4시점 태양위치 × ray-cast",
              "4. shadow union / 4 → [0,1]",
              "5. 격자 centroid 픽셀 매칭",
              "",
              "산출:",
              "heukseok_shadow_accum.tif",
              "grid_heukseok_natural.csv"]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = t
        _set_font(r, size=11, color=COLOR_DARK)
        p.space_before = Pt(1)
    _add_footer(s, 9)

    # ────── 10. OSMnx 교차로 + 횡단보도 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "OSMnx · 교차로 + 횡단보도",
                    "보행자가 모이는 결집 지점을 보행로 옆 추천 가중치로")
    _add_bullets(s, Inches(0.75), Inches(1.8), Inches(7.0), Inches(4.5), [
        "OSMnx — OpenStreetMap Overpass API Python 래퍼",
        "walkable highway 필터: primary~footway·pedestrian 13종",
        "교차로 노드: street_count ≥ 3 (T자·X자)",
        "횡단보도: 한국 OSM은 footway=crossing 패턴 우세",
        "  → highway·footway·crossing 3종 태그 통합 검색",
        "결집지점 union → 격자별 반경 80m 내 개수 + 50m 근접 필터",
        "결과: 동작구 격자 평균 2.02개 결집지점, 노량진 TOP10 진입",
    ], size=13)

    box = _box(s, Inches(8.0), Inches(1.8), Inches(4.8), Inches(4.5),
                RGBColor(0xE0, 0xF7, 0xFA), COLOR_CYAN)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "📍 동작구 통계"
    _set_font(r, size=13, bold=True, color=COLOR_CYAN)
    for t in ["",
              "walkable edges: 수천 개",
              "교차로 노드: 약 1,400 (full 그래프)",
              "횡단보도: 306개 (3종 태그)",
              "  · highway=crossing: 179",
              "  · footway=crossing: 126",
              "  · 중복 제거",
              "",
              "흑석동 안:",
              "  · 교차로 63 · 횡단보도 4"]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = t
        _set_font(r, size=11, color=COLOR_DARK)
        p.space_before = Pt(1)
    _add_footer(s, 10)

    # ────── 11. 시나리오 6개 + 강건 입지 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "시나리오 민감도 6개 + 강건 입지",
                    "정책 관점별 가중치 × 동일 피처셋 → TOP10 변화 추적")
    rows = [
        ("시나리오", "강조 가중치", "최고 Score", "유니크"),
        ("기본", "균형 (18/18/18)", "0.551", "0곳"),
        ("고령자 중시", "vuln 0.38", "0.518", "0곳"),
        ("폭염 중시", "lst 0.38", "0.607", "1곳"),
        ("유동인구 중시", "pop 0.38", "0.648", "3곳"),
        ("보행환경 중시", "streetview 0.32", "0.589", "0곳"),
        ("교차로 중시 (NEW)", "intersection 0.40", "0.474", "0곳"),
    ]
    widths = [Inches(3.3), Inches(3.3), Inches(2.0), Inches(2.0)]
    _table(s, Inches(0.75), Inches(1.8), sum(widths, Inches(0)),
            Inches(2.5), rows, widths, header_size=13, body_size=12,
            bold_rows=[3, 4, 5, 6])

    box = _box(s, Inches(0.75), Inches(4.6), Inches(12), Inches(2.0),
                RGBColor(0xFA, 0xFA, 0xFA), COLOR_ACCENT, 2.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "⭐ 강건 입지 — 6개 시나리오 모두 공통 추천"
    _set_font(r, size=15, bold=True, color=COLOR_ACCENT)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = "동작대로 사당-이수 축 #1 (37.4907, 126.9647) · #2 (37.4898, 126.9670)"
    _set_font(r, size=17, bold=True, color=COLOR_DARK)
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run()
    r.text = "→ 예산 제약 시 최우선 설치 대상 (정책 관점 불변)"
    _set_font(r, size=12, color=COLOR_SUB)
    p3.space_before = Pt(6)
    _add_footer(s, 11)

    # ────── 12. 1차 발표 3대 구역 검증 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "1차 발표 3대 집중구역 정량 검증",
                    "사전 선정한 후보 구역을 알고리즘이 얼마나 재현했는가?")
    rows = [
        ("구역", "판정", "반경 내 TOP", "최근접", "해석"),
        ("노량진역-수산시장", "✓", "2곳 (TOP6, TOP9)", "159m",
         "횡단보도 통합으로 신규 진입"),
        ("사당역-이수역 축", "✓", "4곳 (TOP1·2·4·5)", "253m",
         "환승 대기 + 오후 열스트레스 최상위"),
        ("상도로 주거축", "△", "0곳", "1,516m",
         "보행로 5m 강화 후 후보 줄어듦"),
    ]
    widths = [Inches(2.8), Inches(0.9), Inches(2.7), Inches(1.4), Inches(4.7)]
    _table(s, Inches(0.4), Inches(1.8), sum(widths, Inches(0)),
            Inches(2.5), rows, widths, header_size=13, body_size=12)

    _add_text(s, Inches(0.75), Inches(4.8), Inches(12), Inches(2.0),
              "💡 노량진역-수산시장이 처음으로 TOP10 진입 — OSMnx 횡단보도(footway=crossing) 다중 태그 검색의 직접적 효과.\n"
              "사당-이수 환승 축은 모든 변화에도 흔들리지 않는 강건 입지 #1·#2 차지.\n"
              "상도로는 1차 발표 때 ✓였으나 보행로 필터 5m 강화로 컷 — 정책적 재검토 필요.",
              size=13, color=COLOR_DARK)
    _add_footer(s, 12)

    # ────── 13. 흑석동 — 실측 vs 신규 추천 (핵심 임팩트) ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "흑석동 — 실측 그늘막 18개 vs 신규 추천 TOP 10",
                    "사용자 제공 실데이터로 사각지대·보강 후보 정량 비교")

    rows = [
        ("Rank", "score", "기존그늘막", "분류"),
        ("TOP1", "+0.276", "45.7m", "● 보강"),
        ("TOP2", "+0.231", "44.4m", "● 보강"),
        ("TOP3", "+0.134", "322.9m", "⭐ 사각지대"),
        ("TOP4", "+0.107", "359.6m", "⭐ 사각지대 (최대)"),
        ("TOP5", "+0.100", "234.8m", "⭐ 사각지대"),
        ("TOP6", "+0.072", "98.5m", "○ 중간"),
        ("TOP7", "+0.068", "51.0m", "● 보강"),
        ("TOP8", "+0.057", "153.8m", "⭐ 사각지대"),
        ("TOP9", "+0.016", "70.0m", "○ 중간"),
        ("TOP10", "−0.008", "29.2m", "● 보강"),
    ]
    widths = [Inches(0.9), Inches(1.3), Inches(1.5), Inches(2.5)]
    _table(s, Inches(0.4), Inches(1.7), sum(widths, Inches(0)),
            Inches(4.5), rows, widths, header_size=12, body_size=11)

    # 우측 인사이트 박스
    box = _box(s, Inches(6.7), Inches(1.7), Inches(6.2), Inches(2.2),
                RGBColor(0xFF, 0xF4, 0xE5), COLOR_ACCENT, 2)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "⭐ 사각지대 4곳 — 정책 어필 포인트"
    _set_font(r, size=14, bold=True, color=COLOR_ACCENT)
    for t in [
        "TOP4 (37.49786, 126.96129) — 최대 359.6m",
        "TOP3 (37.49874, 126.95789) — 322.9m",
        "TOP5 (37.51138, 126.96120) — 234.8m",
        "TOP8 (37.50328, 126.96465) — 153.8m",
    ]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = "• " + t
        _set_font(r, size=12, color=COLOR_DARK)
        p.space_before = Pt(3)

    box = _box(s, Inches(6.7), Inches(4.1), Inches(6.2), Inches(2.0),
                RGBColor(0xF5, 0xF8, 0xFD), COLOR_PRIMARY, 2)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "📊 흑석동 데이터 자체"
    _set_font(r, size=14, bold=True, color=COLOR_PRIMARY)
    for t in [
        "실측 그늘막 18개 (고정형 13 + 스마트형 5)",
        "NGII 건물 3,775동 + DSM 4시점 누적 그림자",
        "그늘막이 흑석동 75% 영역 커버 — 정책: 다른 동 우선",
    ]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = "• " + t
        _set_font(r, size=12, color=COLOR_DARK)
        p.space_before = Pt(3)
    _add_footer(s, 13)

    # ────── 14. 예산 제약 최적화 (NEW) ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "예산 제약 최적화 (Budget-Constrained Knapsack)",
                    "예산 N원 = 그늘막 K개 = 어디? 정책 결정에 직접 답")
    box = _box(s, Inches(0.75), Inches(1.8), Inches(7.5), Inches(2.0),
                COLOR_DARK)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "max Σ score_i × x_i"
    _set_font(r, size=18, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    for t in [
        "subject to:",
        "  Σ cost_i × x_i ≤ 예산",
        "  x_i + x_j ≤ 1   if dist(i,j) < 200m",
        "  x_i ∈ {0, 1}",
    ]:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run(); r.text = t
        _set_font(r, size=12, color=RGBColor(0xFF, 0xFF, 0xFF))
        p2.space_before = Pt(2)

    _add_bullets(s, Inches(0.75), Inches(4.0), Inches(7.5), Inches(2.5), [
        "PuLP 정수 선형계획법 (CBC 솔버)",
        "후보 풀: 필터 통과 격자 score 상위 50개",
        "예산: 4,000만원 / 그늘막 단가: 800만원 → 최대 5개",
        "공간 분산 제약: 선정 격자 간 ≥ 200m 이격 (29쌍 제약)",
        "결과: status=Optimal, 5개 선정, 총 score 1.852, 예산 100% 사용",
    ], size=13)

    box2 = _box(s, Inches(8.5), Inches(1.8), Inches(4.4), Inches(4.5),
                 RGBColor(0xFF, 0xF4, 0xE5), COLOR_ACCENT, 2)
    tf = box2.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "🎯 선정 5곳 (분산 균형)"
    _set_font(r, size=13, bold=True, color=COLOR_ACCENT)
    for t in [
        "",
        "#1 (37.490, 126.964) s=0.466",
        "    이수역 인근",
        "#2 (37.493, 126.962) s=0.363",
        "    이수 북쪽",
        "#3 (37.491, 126.968) s=0.359",
        "    사당역 동쪽",
        "#4 (37.486, 126.967) s=0.333",
        "    사당역 남쪽",
        "#5 (37.513, 126.944) s=0.331",
        "    노량진역 인근",
    ]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = t
        _set_font(r, size=11, color=COLOR_DARK)
        p.space_before = Pt(1)
    _add_footer(s, 14)

    # ────── 15. Streamlit 대시보드 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "인터랙티브 대시보드 (Streamlit)",
                    "가중치 슬라이더 → TOP 10 실시간 재계산")
    _add_bullets(s, Inches(0.75), Inches(1.8), Inches(6.5), Inches(4.5), [
        "`streamlit run app.py` → http://localhost:8501",
        "사이드바: 시나리오 프리셋 + 7개 가중치 슬라이더",
        "본문: Folium 지도 임베드 + TOP 10 표 + 통계",
        "캐싱: @st.cache_resource 로 격자·피처 1회 계산,",
        "  슬라이더 조작은 가중합 재계산만 (밀리초 반응)",
        "발표 현장 \"vuln 더 올리면?\" 질문에 즉답",
    ], size=14)
    _placeholder_image(s, Inches(7.5), Inches(1.8),
                       Inches(5.3), Inches(4.5),
                       "Streamlit 대시보드 스크린샷")
    _add_footer(s, 15)

    # ────── 16. 한계 + 한 줄 요약 ──────
    s = prs.slides.add_slide(blank)
    _add_header_bar(s, "한계 · 다음 단계 · 한 줄 요약", None)

    _add_text(s, Inches(0.75), Inches(1.7), Inches(6), Inches(0.4),
              "현재 한계", size=16, bold=True, color=COLOR_ACCENT)
    _add_bullets(s, Inches(0.75), Inches(2.2), Inches(6.5), Inches(2.8), [
        "5종 기본 데이터 중 4종 더미 (파이프라인은 자동 전환)",
        "LST 래스터 직접 파싱 미연결 (Landsat ST_B10)",
        "SAM 추출 건물에 추정 높이 부여 (NGII는 실측)",
        "예산·단가는 더미 800만원/개 (실 단가 확보 시 교체)",
    ], size=13)

    _add_text(s, Inches(7.5), Inches(1.7), Inches(6), Inches(0.4),
              "다음 단계", size=16, bold=True, color=COLOR_PRIMARY)
    _add_bullets(s, Inches(7.5), Inches(2.2), Inches(5.5), Inches(2.8), [
        "서울 열린데이터광장 API 연결 (P0)",
        "동작구 전 동 실측 그늘막 (흑석동만 18개 확보)",
        "Streamlit 슬라이더에 시각화 다양화",
        "타 자치구 확장 (BBOX 교체로 재사용)",
    ], size=13)

    box = _box(s, Inches(0.75), Inches(5.4), Inches(12), Inches(1.4),
                COLOR_DARK)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "항공사진·거리뷰·DSM·도로망·실측 그늘막을 컴퓨터비전·GIS로 통합"
    _set_font(r, size=17, bold=True, color=RGBColor(0xFF, 0xEE, 0x58))
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = "→ 6 시나리오 모두에서 살아남는 강건 입지 2곳 + 흑석동 사각지대 4곳 도출"
    _set_font(r, size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    p2.space_before = Pt(6)
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run()
    r.text = "SAM + SegFormer + Deep Umbra(DSM) + OSMnx + NGII + 실측 그늘막 + MCDA + Knapsack"
    _set_font(r, size=11, color=RGBColor(0xCC, 0xCC, 0xCC))
    p3.space_before = Pt(4)
    _add_footer(s, 16)

    # 각 슬라이드에 스피커 노트 삽입 (PPTX notes_slide)
    for i, slide in enumerate(prs.slides, 1):
        note_text = SLIDE_NOTES.get(i, "")
        if note_text:
            tf = slide.notes_slide.notes_text_frame
            tf.text = note_text.strip()

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"[완료] {path}")
    print(f"  슬라이드 {TOTAL}장 · 16:9 · 맑은 고딕")
    print(f"  각 슬라이드에 스피커 노트 자동 삽입 (PPTX 노트 영역)")
    print(f"  상세 노트: presentation/speaker_notes.md")
